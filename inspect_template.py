
from pptx import Presentation
import os # Added import for os module
import sys # Added import for sys module for exit()

TEMPLATE_PATH = "d:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx" # Renamed and capitalized template_path

with open("template_info_v3.txt", "w", encoding="utf-8") as f:
    f.write(f"--- Inspecting Template: {str(TEMPLATE_PATH)} ---\n")
    if not os.path.exists(TEMPLATE_PATH):
        f.write(f"Error: Template not found at {TEMPLATE_PATH}\n")
        sys.exit(1) # Changed exit() to sys.exit() for clarity and best practice

    prs = Presentation(TEMPLATE_PATH)

    f.write("\n--- Existing Slides ---\n") # Changed "Existing Slides Count" to "Existing Slides"
    for i, slide in enumerate(prs.slides):
        f.write(f"Slide {i}: Layout '{slide.slide_layout.name}'\n")
        for shape in slide.placeholders: # Changed from slide.shapes to slide.placeholders
             f.write(f"  Placeholder idx={shape.placeholder_format.idx}, name='{shape.name}', type={shape.placeholder_format.type}\n")
             if hasattr(shape, "text_frame") and shape.text_frame.text: # Added check for text_frame and its content
                f.write(f"    - Text: {shape.text_frame.text[:50]}...\n")

    f.write("\n--- Available Layouts ---\n")
    f.write(f"Total Layouts: {len(prs.slide_layouts)}\n")
    for i, layout in enumerate(prs.slide_layouts):
        if i > 15: break # Changed break condition from 10 to 15
        f.write(f"\nLayout {i}: {layout.name}\n")
        for ph in layout.placeholders:
            f.write(f"  Placeholder idx={ph.placeholder_format.idx}, name='{ph.name}', type={ph.placeholder_format.type} ({ph.placeholder_format.type})\n")
