"""
Evaluation Runner for AI_Pedia_Local_stream System.

This script automates the evaluation of the multi-agent content generation system.
It runs the system for a given list of topics, generates artifacts (PPTX, Video, Code, Quiz),
and computes quality, consistency, and system-level metrics.

Usage:
    python tests/evaluation_runner.py --topics "Machine Learning Basics" "Python for Data Science" --output report.csv
    python tests/evaluation_runner.py --file topics.json --output report.csv
"""

import argparse
import asyncio
import csv
import json
import logging
import os
import re
import subprocess
import sys
import time
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

# Add project root to sys.path
# Since we are running from project root (d:\L3\Individual_project\AI_Pedia_Local_stream),
# PROJECT_ROOT should resolve to that.
PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

# Also add current directory if running from there
cwd = os.getcwd()
if cwd not in sys.path:
    sys.path.insert(0, cwd)

# Import system components
try:
    # Try absolute import first given the structure
    from manager_agent.task_manager_agent import stream_workflow
    from config import GENERATED_DIR
except ImportError as e:
    print(f"Import Error: {e}")
    # Fallback: try to find where we are
    print(f"Current sys.path: {sys.path}")
    print(f"Current CWD: {os.getcwd()}")
    sys.exit(1)

# Import evaluation libraries
try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed. Install via: pip install python-pptx")
    import types
    # Mock for development if needed, but better to fail early
    Presentation = None

# Import specific modules (Assuming they are in tests/)
try:
    from tests.observability_module import TaskContext, LLMWrapper
    from tests.learning_gain_evaluator import LearningSimulator, KnowledgeTracer
except ImportError:
    # Use relative import style if running as module, or adjust path
    # Since PROJECT_ROOT is added, this should work
    from observability_module import TaskContext, LLMWrapper
    from learning_gain_evaluator import LearningSimulator, KnowledgeTracer

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("evaluation.log", mode="w", encoding="utf-8"),
    ],
)
logger = logging.getLogger(__name__)

# --- Configuration ---
MIN_VIDEO_SIZE_BYTES = 100 * 1024  # 100 KB
CODE_EXEC_TIMEOUT_SECONDS = 600  # 10 minutes

class Evaluator:
    def __init__(self, output_csv: str = "evaluation_report.csv"):
        self.output_csv = output_csv
        self.results = []

    async def run_evaluation(self, topics: List[str]):
        """Main evaluation loop."""
        logger.info(f"Starting evaluation for {len(topics)} topics.")
        
        for topic in topics:
            logger.info(f"=== Processing Topic: {topic} ===")
            metrics = await self.process_topic(topic)
            self.results.append(metrics)
            logger.info(f"Finished topic: {topic}. Success: {metrics.get('success_flag', False)}")

        self.write_report()
        logger.info(f"Evaluation complete. Report saved to {self.output_csv}")

    async def process_topic(self, topic: str) -> Dict[str, Any]:
        """Runs the system for a single topic and evaluates outputs."""
        start_time = time.time()
        
        # 1. Run System
        artifacts = {
            "presentation": [],
            "video": [],
            "code": [],
            "quiz": [] 
        }
        video_script = ""  # New: Capture script
        
        # Configuration to force all agents
        config = {
            "video": True,
            "slides": True,
            "code": True,
            "quizzes": True
        }
        
        run_success = True
        try:
            # stream_workflow is a sync generator
            iterator = stream_workflow(topic, config)
            
            for event_str in iterator:
                # event_str format: "event: type\ndata: json\n\n"
                # Parse event manually
                lines = event_str.strip().split("\n")
                event_type = ""
                data_str = ""
                
                for line in lines:
                    if line.startswith("event: "):
                        event_type = line[len("event: "):].strip()
                    elif line.startswith("data: "):
                        data_str = line[len("data: "):].strip()
                
                if not data_str:
                    continue
                    
                try:
                    data = json.loads(data_str)
                except json.JSONDecodeError:
                    continue

                if event_type == "artifact":
                    art_type = data.get("type", "file")
                    url = data.get("url", "")
                    # Map to our categories
                    if art_type == "presentation" or url.endswith(".pptx"):
                        artifacts["presentation"].append(url)
                    elif art_type == "video" or url.endswith(".mp4"):
                        artifacts["video"].append(url)
                    elif art_type == "code" or url.endswith(".py"):
                        artifacts["code"].append(url)
                
                elif event_type == "quiz":
                    # Quiz event usually contains the content in data['content']
                    artifacts["quiz"].append(data)
                    
                elif event_type == "script":  # New Event Handler
                    video_script = data.get("content", "")
                
                elif event_type == "error":
                    logger.error(f"System Error: {data.get('content')}")
                    run_success = False

        except Exception as e:
            logger.error(f"Error running workflow: {e}")
            run_success = False

        end_time = time.time()
        latency = end_time - start_time
        
        # 2. Evaluate Artifacts (Basic Quality)
        slide_metrics = self.evaluate_slide_quality(artifacts["presentation"][0]) if artifacts["presentation"] else {}
        video_metrics = self.evaluate_video_quality(artifacts["video"][0]) if artifacts["video"] else {}
        code_metrics = self.evaluate_code_quality(artifacts["code"][0], topic) if artifacts["code"] else {}
        quiz_metrics = self.evaluate_quiz_quality(artifacts["quiz"])
        
        # 3. Consistency Check (Keyword Based)
        consistency_metrics = self.evaluate_consistency(
            slide_metrics.get("extracted_text", ""),
            code_metrics.get("extracted_text", ""),
            topic
        )
        
        # 4. Learning Gain Evaluation (Deep Semantic Check)
        learning_metrics = {
            "accuracy_pre": 0.0,
            "accuracy_post": 0.0,
            "learning_gain": 0.0,
            "coverage_rate": 0.0,
            "eval_latency": 0.0
        }
        
        # Only run if we have a valid quiz and at least some content
        slides_text = slide_metrics.get("extracted_text", "")
        code_text = code_metrics.get("extracted_text", "")
        
        if artifacts["quiz"] and (slides_text or video_script or code_text):
            logger.info("Starting Learning Gain Evaluation...")
            eval_start = time.time()
            try:
                # Extract quiz questions from the event data
                # content might be wrapped
                quiz_data = artifacts["quiz"][0].get("content")
                if isinstance(quiz_data, str):
                    try: quiz_data = json.loads(quiz_data)
                    except: quiz_data = []
                
                if not isinstance(quiz_data, list):
                     # If it's a dict wrapper like {"questions": [...]}, extract
                     if isinstance(quiz_data, dict) and "questions" in quiz_data:
                         quiz_data = quiz_data["questions"]
                     else:
                         quiz_data = []

                if quiz_data:
                    safe_topic_name = "".join(x for x in topic if x.isalnum() or x in " _-").strip().replace(" ", "_")
                    run_dir = Path("evaluation_runs") / f"{safe_topic_name}_{int(time.time())}"
                    run_dir.mkdir(parents=True, exist_ok=True)
                    
                    # Initialize Observability & Evaluator
                    with TaskContext(task_id=f"eval_{safe_topic_name}", run_dir=run_dir) as ctx:
                        llm = LLMWrapper(ctx)
                        
                        # --- Mocking for Demo/Dev environment validity ---
                        # In production, remove this to use real API keys
                        if not os.getenv("OPENAI_API_KEY"):
                             # If no key, inject mock to prevent crash during demo
                            original_chat = llm.chat
                            def mock_chat_override(agent_name, model_name, messages, **kwargs):
                                if agent_name.startswith("student"):
                                    return {
                                        "choices": [{"message": {"content": json.dumps({"answer": "A", "confidence": 0.6})}}],
                                        "usage": {"total_tokens": 50}
                                    }
                                elif agent_name == "knowledge_tracer":
                                    return {
                                        "choices": [{"message": {"content": json.dumps({"covered": True, "reason": "Found"})}}],
                                        "usage": {"total_tokens": 50}
                                    }
                                return original_chat(agent_name, model_name, messages, **kwargs)
                            llm.chat = mock_chat_override
                        # -------------------------------------------------

                        sim = LearningSimulator(llm)
                        tracer = KnowledgeTracer(llm)

                        # A. Pre-test
                        pre = sim.run_pretest(quiz_data)

                        # B. Study
                        sim.run_study(slides_text, video_script, code_text)

                        # C. Post-test
                        combined_materials = f"Slides:\n{slides_text}\n\nVideo Script:\n{video_script}\n\nCode:\n{code_text}"
                        post = sim.run_posttest(quiz_data, combined_materials)

                        # D. Coverage
                        cov = tracer.evaluate_coverage(quiz_data, combined_materials)

                        learning_metrics["accuracy_pre"] = pre["accuracy"]
                        learning_metrics["accuracy_post"] = post["accuracy"]
                        learning_metrics["learning_gain"] = post["accuracy"] - pre["accuracy"]
                        learning_metrics["coverage_rate"] = cov["coverage_rate"]
                        
            except Exception as e:
                logger.error(f"Learning Evaluation Failed: {e}")
                
            learning_metrics["eval_latency"] = time.time() - eval_start

        # 5. Compile Final Metrics
        final_metrics = {
            "topic": topic,
            "timestamp": datetime.now().isoformat(),
            "latency_seconds": round(latency, 2),
            "success_flag": run_success,
            
            # Artifact Metrics
            "slide_count": slide_metrics.get("slide_count", 0),
            "video_pass_size": video_metrics.get("pass_size", False),
            "code_exec_success": code_metrics.get("exec_success", False),
            "quiz_count": quiz_metrics.get("count", 0),
            
            # Learning & Quality Metrics
            "consistency_score": float(consistency_metrics.get("score", 0.0)),
            "accuracy_pre": float(learning_metrics["accuracy_pre"]),
            "accuracy_post": float(learning_metrics["accuracy_post"]),
            "learning_gain": float(learning_metrics["learning_gain"]),
            "coverage_rate": float(learning_metrics["coverage_rate"]),
            
            # System
            "token_usage": "See Trace Logs"  # Handled by observability module
        }
        
        return final_metrics

    def evaluate_slide_quality(self, ppt_path_str: str) -> Dict[str, Any]:
        """Evaluates PPTX file structure and content."""
        metrics = {"slide_count": 0, "has_summary_slide": False, "keyword_score": 0.0, "extracted_text": ""}
        path = Path(ppt_path_str)
        
        if not path.exists():
            logger.warning(f"PPTX file not found: {path}")
            return metrics

        try:
            if Presentation is None:
                logger.warning("python-pptx not available, skipping detailed slide checks.")
                return metrics
                
            prs = Presentation(path)
            metrics["slide_count"] = len(prs.slides)
            
            all_text = []
            titles = []
            
            for slide in prs.slides:
                if slide.shapes.title and slide.shapes.title.text:
                    titles.append(slide.shapes.title.text.lower())
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
            
            text_content = " ".join(all_text).lower()
            metrics["extracted_text"] = text_content
            
            # Rule-based checks
            if any("summary" in t or "conclusion" in t for t in titles):
                metrics["has_summary_slide"] = True
                
        except Exception as e:
            logger.error(f"Error reading PPTX: {e}")

        return metrics

    def evaluate_code_quality(self, code_path_str: str, topic: str) -> Dict[str, Any]:
        """Evaluates Python code syntax and execution."""
        metrics = {"exec_success": False, "import_checks": False, "syntax_valid": False, "extracted_text": ""}
        path = Path(code_path_str)
        
        if not path.exists():
            return metrics

        try:
            code_content = path.read_text(encoding="utf-8")
            metrics["extracted_text"] = code_content
            
            # Syntax Check
            compile(code_content, str(path), 'exec')
            metrics["syntax_valid"] = True
            
            # Import Checks (Basic heuristic)
            if "machine learning" in topic.lower() or "data" in topic.lower():
                if "import pandas" in code_content or "import numpy" in code_content or "import sklearn" in code_content or "import torch" in code_content:
                    metrics["import_checks"] = True
            else:
                # pass by default if no specific domain logic
                metrics["import_checks"] = True

            # Execution Check
            # Run in a separate process with timeout
            metrics["exec_success"] = False # Default unless proven true
            # Safety: Basic check for destructive commands
            if "os.system" in code_content or "subprocess" in code_content or "shutil.rmtree" in code_content:
                logger.warning(f"Skipping execution of {path} due to potential unsafe operations.")
            else:
                try:
                    # Run within a constrained environment if possible, but here using subprocess
                    # Run in file's directory so it finds assets
                    cwd = path.parent
                    
                    result = subprocess.run(
                        [sys.executable, str(path)],
                        capture_output=True,
                        text=True,
                        timeout=CODE_EXEC_TIMEOUT_SECONDS,
                        cwd=cwd 
                    )
                    
                    if result.returncode == 0:
                        metrics["exec_success"] = True
                    else:
                        logger.warning(f"Code execution failed: {result.stderr[:200]}")
                        
                except subprocess.TimeoutExpired:
                    pass
                except Exception:
                    pass

        except SyntaxError:
            metrics["syntax_valid"] = False
        except Exception as e:
            logger.error(f"Error evaluating code: {e}")

        return metrics

    def evaluate_video_quality(self, video_path_str: str) -> Dict[str, Any]:
        """Evaluates Video file existence and size."""
        metrics = {"exists": False, "pass_size": False}
        path = Path(video_path_str)
        
        if path.exists():
            metrics["exists"] = True
            try:
                size = path.stat().st_size
                if size >= MIN_VIDEO_SIZE_BYTES:
                    metrics["pass_size"] = True
            except Exception:
                pass
        
        return metrics

    def evaluate_quiz_quality(self, quizzes: List[Dict]) -> Dict[str, Any]:
        """Evaluates generated quizzes."""
        metrics = {"count": 0, "valid_json": True} # events are already parsed JSON
        
        if not quizzes:
            metrics["valid_json"] = False
            return metrics
            
        metrics["count"] = len(quizzes)
        return metrics

    def evaluate_consistency(self, slide_text: str, code_text: str, topic: str) -> Dict[str, Any]:
        """Computes keyword consistency between slides and code."""
        metrics = {"score": 0.0}
        
        # Generate keywords from topic
        # Simple heuristic: split topic into words, remove stop words
        stop_words = {"a", "an", "the", "in", "on", "at", "for", "to", "of", "and", "or", "with", "by", "is", "basics", "introduction", "tutorial"}
        keywords = {w.lower() for w in topic.split() if w.lower() not in stop_words and len(w) > 2}
        
        if not keywords:
            metrics["score"] = 1.0 # Trivial case
            return metrics
            
        slide_text_lower = slide_text.lower()
        code_text_lower = code_text.lower()
        
        matches = 0
        total_checks = 0
        
        for kw in keywords:
            in_slides = kw in slide_text_lower
            in_code = kw in code_text_lower
            
            # We want both ideally
            if in_slides and in_code:
                matches += 1
            elif in_slides or in_code:
                matches += 0.5 # Partial credit
            
            total_checks += 1
                
        metrics["score"] = matches / total_checks if total_checks > 0 else 0.0
        return metrics

    def write_report(self):
        """Writes the evaluation results to CSV."""
        if not self.results:
            logger.warning("No results to write.")
            return

        headers = list(self.results[0].keys())
        
        try:
            with open(self.output_csv, "w", newline="", encoding="utf-8") as f:
                writer = csv.DictWriter(f, fieldnames=headers)
                writer.writeheader()
                writer.writerows(self.results)
        except Exception as e:
            logger.error(f"Error writing CSV report: {e}")

def main():
    parser = argparse.ArgumentParser(description="Run system evaluation.")
    parser.add_argument("--topics", nargs="+", help="List of topics to evaluate") # e.g. python tests/evaluation_runner.py --topics "Machine Learning" "Deep Learning"
    parser.add_argument("--file", help="JSON file containing topics") 
    parser.add_argument("--output", default="evaluation_report.csv", help="Output CSV path")
    
    args = parser.parse_args()
    
    topics = []
    if args.topics:
        topics = args.topics
    elif args.file:
        try:
            with open(args.file, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, list):
                    topics = data
                elif isinstance(data, dict) and "topics" in data:
                    topics = data["topics"]
        except Exception as e:
            print(f"Error reading topics file: {e}")
            sys.exit(1)
            
    if not topics:
        print("No topics provided. Use --topics 'Topic 1' 'Topic 2' or --file topics.json")
        sys.exit(1)

    evaluator = Evaluator(output_csv=args.output)
    
    try:
        asyncio.run(evaluator.run_evaluation(topics))
    except KeyboardInterrupt:
        print("\nEvaluation interrupted by user.")
    except Exception as e:
        print(f"\nCritical Error: {e}")

if __name__ == "__main__":
    main()
