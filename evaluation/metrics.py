"""Deterministic evaluation metrics for AI-Pedia outputs."""

from __future__ import annotations

import json
import logging
import math
import re
import subprocess
import sys
import time
from collections import Counter
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

from evaluation.constants import (
    CODE_EXEC_TIMEOUT_SEC,
    EXPECTED_QUIZ_OPTION_COUNT,
    EXPECTED_QUIZ_QUESTION_COUNT,
    PIP_INSTALL_TIMEOUT_SEC,
    TFIDF_TOP_K,
    VALID_ANSWER_LETTERS,
)

logger = logging.getLogger(__name__)

_SUMMARY_KEYWORDS = {
    "summary",
    "conclusion",
    "recap",
    "wrap up",
    "wrap-up",
    "key takeaways",
    "takeaway",
    "final thoughts",
}

_SLIDE_TITLE_LAYOUT_HINTS = {"title slide", "title"}

_STOPWORDS = {
    "a", "about", "above", "after", "again", "against", "all", "am", "an", "and", "any", "are",
    "as", "at", "be", "because", "been", "before", "being", "below", "between", "both", "but", "by",
    "can", "did", "do", "does", "doing", "down", "during", "each", "few", "for", "from", "further",
    "had", "has", "have", "having", "he", "her", "here", "hers", "herself", "him", "himself", "his",
    "how", "i", "if", "in", "into", "is", "it", "its", "itself", "just", "me", "more", "most", "my",
    "myself", "no", "nor", "not", "now", "of", "off", "on", "once", "only", "or", "other", "our",
    "ours", "ourselves", "out", "over", "own", "same", "she", "should", "so", "some", "such", "than",
    "that", "the", "their", "theirs", "them", "themselves", "then", "there", "these", "they", "this",
    "those", "through", "to", "too", "under", "until", "up", "very", "was", "we", "were", "what",
    "when", "where", "which", "while", "who", "whom", "why", "will", "with", "you", "your", "yours",
    "yourself", "yourselves",
}

_IRREGULAR_LEMMA = {
    "indices": "index",
    "matrices": "matrix",
    "children": "child",
    "teeth": "tooth",
    "men": "man",
    "women": "woman",
    "data": "data",
}


# ---------------------------------------------------------------------------
# Generic helpers
# ---------------------------------------------------------------------------

def read_text_file(path: str) -> str:
    p = Path(path)
    if not path or not p.is_file():
        return ""
    try:
        return p.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:  # noqa: BLE001
        logger.error("Error reading %s: %s", path, exc)
        return ""


def _clean_whitespace(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "")).strip()


def _jaccard(a: Iterable[str], b: Iterable[str]) -> float:
    sa = set(a)
    sb = set(b)
    if not sa and not sb:
        return 1.0
    if not sa or not sb:
        return 0.0
    return len(sa & sb) / len(sa | sb)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def _extract_pptx_text_data(pptx_path: str) -> Dict[str, Any]:
    out = {
        "slide_count": 0,
        "all_text": "",
        "slide_texts": [],
        "slide_titles": [],
        "slide_layout_names": [],
    }

    path = Path(pptx_path)
    if not pptx_path or not path.is_file():
        return out

    try:
        from pptx import Presentation
    except Exception as exc:  # noqa: BLE001
        logger.warning("python-pptx unavailable, cannot parse slides: %s", exc)
        return out

    try:
        prs = Presentation(str(path))
        all_chunks: List[str] = []
        slide_texts: List[str] = []
        slide_titles: List[str] = []
        layout_names: List[str] = []

        for slide in prs.slides:
            texts: List[str] = []
            for shape in slide.shapes:
                txt = ""
                if hasattr(shape, "text"):
                    txt = _clean_whitespace(str(getattr(shape, "text", "")))
                if txt:
                    texts.append(txt)

            title_text = ""
            if getattr(slide.shapes, "title", None) is not None and slide.shapes.title is not None:
                title_text = _clean_whitespace(str(getattr(slide.shapes.title, "text", "")))

            slide_titles.append(title_text)
            layout_names.append(_clean_whitespace(str(getattr(slide.slide_layout, "name", ""))).lower())

            merged = "\n".join(texts)
            slide_texts.append(merged)
            if merged:
                all_chunks.append(merged)

        out["slide_count"] = len(slide_texts)
        out["all_text"] = "\n\n".join(all_chunks)
        out["slide_texts"] = slide_texts
        out["slide_titles"] = slide_titles
        out["slide_layout_names"] = layout_names
    except Exception as exc:  # noqa: BLE001
        logger.error("Failed to parse pptx %s: %s", pptx_path, exc)

    return out


def extract_text_from_pptx(pptx_path: str) -> str:
    return _extract_pptx_text_data(pptx_path).get("all_text", "")


def extract_slide_texts_from_pptx(pptx_path: str) -> List[str]:
    return list(_extract_pptx_text_data(pptx_path).get("slide_texts", []))


def _classify_slide_type(index: int, title: str, text: str, layout_name: str) -> str:
    title_l = title.lower().strip()
    text_l = text.lower().strip()

    # Summary by title keyword
    if title_l and any(k in title_l for k in _SUMMARY_KEYWORDS):
        return "summary"

    words = re.findall(r"[a-zA-Z0-9]+", text_l)
    word_count = len(words)

    # Title-slide heuristics
    if layout_name in _SLIDE_TITLE_LAYOUT_HINTS or "title slide" in layout_name:
        return "title"

    if index == 0 and title_l and word_count <= 35:
        return "title"

    if title_l and word_count <= 18 and not any(k in title_l for k in _SUMMARY_KEYWORDS):
        return "title"

    return "content"


def evaluate_slides(pptx_path: str) -> Dict[str, Any]:
    data = _extract_pptx_text_data(pptx_path)

    title_count = 0
    content_count = 0
    summary_count = 0
    slide_types: List[str] = []

    for i, text in enumerate(data["slide_texts"]):
        title = data["slide_titles"][i] if i < len(data["slide_titles"]) else ""
        layout = data["slide_layout_names"][i] if i < len(data["slide_layout_names"]) else ""
        slide_type = _classify_slide_type(i, title, text, layout)
        slide_types.append(slide_type)
        if slide_type == "title":
            title_count += 1
        elif slide_type == "summary":
            summary_count += 1
        else:
            content_count += 1

    compliant = (title_count == 1 and content_count >= 1 and summary_count >= 1)

    return {
        "slide_count": int(data["slide_count"]),
        "title_slide_count": title_count,
        "content_slide_count": content_count,
        "summary_slide_count": summary_count,
        "slide_types": slide_types,
        "slide_structural_compliance": 1.0 if compliant else 0.0,
        "slides_pass": bool(compliant),
        "slides_text": data["all_text"],
        "slide_texts": data["slide_texts"],
        "slide_titles": data["slide_titles"],
    }


# ---------------------------------------------------------------------------
# Code
# ---------------------------------------------------------------------------

def _extract_missing_modules(stderr_text: str) -> List[str]:
    text = stderr_text or ""
    patterns = [
        r"No module named ['\"]([A-Za-z0-9_\.\-]+)['\"]",
        r"ModuleNotFoundError:\s*No module named ['\"]([A-Za-z0-9_\.\-]+)['\"]",
    ]

    modules: List[str] = []
    for pat in patterns:
        modules.extend(re.findall(pat, text))

    cleaned: List[str] = []
    for m in modules:
        pkg = m.strip().split(".")[0]
        if pkg and pkg not in cleaned:
            cleaned.append(pkg)
    return cleaned


def _run_python_script(script_path: Path, cwd: Path, timeout_sec: int) -> Tuple[int, str, str, float]:
    t0 = time.perf_counter()
    proc = subprocess.run(
        [sys.executable, str(script_path)],
        cwd=str(cwd),
        capture_output=True,
        text=True,
        timeout=timeout_sec,
    )
    dt = time.perf_counter() - t0
    return proc.returncode, proc.stdout or "", proc.stderr or "", dt


def _install_dependency(package_name: str, cwd: Path, timeout_sec: int) -> Tuple[bool, str]:
    try:
        proc = subprocess.run(
            [sys.executable, "-m", "pip", "install", package_name],
            cwd=str(cwd),
            capture_output=True,
            text=True,
            timeout=timeout_sec,
        )
        ok = proc.returncode == 0
        log = (proc.stdout or "") + "\n" + (proc.stderr or "")
        return ok, log[-4000:]
    except Exception as exc:  # noqa: BLE001
        return False, str(exc)


def evaluate_code(
    code_path: str,
    timeout_sec: int = CODE_EXEC_TIMEOUT_SEC,
    run_dir: str = "",
    allow_auto_install: bool = True,
    pip_install_timeout_sec: int = PIP_INSTALL_TIMEOUT_SEC,
) -> Dict[str, Any]:
    result: Dict[str, Any] = {
        "code_exec_pass": False,
        "code_runtime_sec": 0.0,
        "code_error_type": None,
        "code_error_msg": None,
        "auto_install_attempted": False,
        "auto_install_succeeded": False,
        "auto_installed_packages": [],
        "auto_install_log": "",
    }

    if not code_path:
        result["code_error_type"] = "FileNotFound"
        result["code_error_msg"] = "Empty code path"
        return result

    script = Path(code_path).expanduser()
    if not script.is_absolute():
        script = (Path.cwd() / script).resolve()

    if not script.is_file():
        result["code_error_type"] = "FileNotFound"
        result["code_error_msg"] = f"{script} does not exist"
        return result

    cwd = Path(run_dir).expanduser().resolve() if run_dir else script.parent.resolve()
    if not cwd.is_dir():
        cwd = script.parent.resolve()

    try:
        rc, stdout, stderr, dt = _run_python_script(script, cwd, timeout_sec)
        result["code_runtime_sec"] = round(dt, 3)

        if rc == 0:
            result["code_exec_pass"] = True
            return result

        err_blob = (stderr or "") + "\n" + (stdout or "")
        result["code_error_type"] = "RuntimeError"
        result["code_error_msg"] = err_blob[-4000:]

        if not allow_auto_install:
            return result

        missing = _extract_missing_modules(err_blob)
        if not missing:
            return result

        result["auto_install_attempted"] = True
        logs: List[str] = []
        installed: List[str] = []

        for pkg in missing:
            ok, log = _install_dependency(pkg, cwd=cwd, timeout_sec=pip_install_timeout_sec)
            logs.append(f"[{pkg}] {'ok' if ok else 'fail'}\n{log}")
            if ok:
                installed.append(pkg)

        result["auto_installed_packages"] = installed
        result["auto_install_log"] = "\n\n".join(logs)[-8000:]

        # Retry execution once after install attempts
        rc2, stdout2, stderr2, dt2 = _run_python_script(script, cwd, timeout_sec)
        result["code_runtime_sec"] = round(result["code_runtime_sec"] + dt2, 3)

        if rc2 == 0:
            result["code_exec_pass"] = True
            result["code_error_type"] = None
            result["code_error_msg"] = None
            result["auto_install_succeeded"] = True
        else:
            result["code_error_type"] = "RuntimeError"
            result["code_error_msg"] = ((stderr2 or "") + "\n" + (stdout2 or ""))[-4000:]

        return result

    except subprocess.TimeoutExpired:
        result["code_runtime_sec"] = float(timeout_sec)
        result["code_error_type"] = "Timeout"
        result["code_error_msg"] = f"Exceeded {timeout_sec} sec"
        return result
    except Exception as exc:  # noqa: BLE001
        result["code_error_type"] = type(exc).__name__
        result["code_error_msg"] = str(exc)[:1000]
        return result


# ---------------------------------------------------------------------------
# Quiz
# ---------------------------------------------------------------------------

def _normalize_answer_letter(raw_answer: Any, options: Sequence[Any]) -> Optional[str]:
    if raw_answer is None:
        return None

    ans = str(raw_answer).strip()
    if not ans:
        return None

    upper = ans.upper()

    if upper in VALID_ANSWER_LETTERS:
        return upper

    if upper == "0":
        return "A"

    if upper in {"1", "2", "3", "4"}:
        return chr(ord("A") + int(upper) - 1)

    # Match answer text to option text
    ans_norm = _clean_whitespace(ans).lower()
    for i, opt in enumerate(options):
        if _clean_whitespace(str(opt)).lower() == ans_norm:
            return chr(ord("A") + i)

    # Match leading letter style: "A) ..."
    m = re.match(r"^([A-Da-d])[\)\.:\-\s].*", ans)
    if m:
        return m.group(1).upper()

    return None


def load_quiz_questions(quiz_json_path: str) -> Dict[str, Any]:
    out = {
        "title": "",
        "questions": [],
        "raw": None,
        "error": None,
    }

    path = Path(quiz_json_path)
    if not quiz_json_path or not path.is_file():
        out["error"] = "Quiz file not found"
        return out

    try:
        raw = json.loads(path.read_text(encoding="utf-8"))
    except Exception as exc:  # noqa: BLE001
        out["error"] = f"Invalid JSON: {exc}"
        return out

    out["raw"] = raw

    if isinstance(raw, dict):
        out["title"] = str(raw.get("title", ""))
        questions = raw.get("questions", raw.get("quiz", []))
    elif isinstance(raw, list):
        questions = raw
    else:
        out["error"] = "Quiz root must be dict or list"
        return out

    normalized: List[Dict[str, Any]] = []
    for i, q in enumerate(questions):
        if not isinstance(q, dict):
            continue

        opts = q.get("options", [])
        if not isinstance(opts, list):
            opts = []
        options = [str(x) for x in opts]

        raw_answer = q.get("answer", q.get("correct_answer"))
        answer_letter = _normalize_answer_letter(raw_answer, options)

        normalized.append(
            {
                "id": q.get("id", i + 1),
                "question": str(q.get("question", "")),
                "options": options,
                "answer_letter": answer_letter,
                "explanation": str(q.get("explanation", "")),
                "type": str(q.get("type", "multiple_choice")),
            }
        )

    out["questions"] = normalized
    return out


def _quiz_text_blob(questions: Sequence[Dict[str, Any]]) -> str:
    chunks: List[str] = []
    for q in questions:
        chunks.append(str(q.get("question", "")))
        for opt in q.get("options", []):
            chunks.append(str(opt))
        chunks.append(str(q.get("explanation", "")))
    return "\n".join([c for c in chunks if c]).strip()


def evaluate_quiz(
    quiz_json_path: str,
    expected_n_questions: int = EXPECTED_QUIZ_QUESTION_COUNT,
    expected_n_options: int = EXPECTED_QUIZ_OPTION_COUNT,
) -> Dict[str, Any]:
    loaded = load_quiz_questions(quiz_json_path)
    questions = loaded["questions"]

    valid_answer_n = 0
    valid_format_n = 0
    has_explanation_n = 0

    for q in questions:
        q_text = _clean_whitespace(str(q.get("question", "")))
        options = q.get("options", [])
        answer_letter = q.get("answer_letter")
        explanation = _clean_whitespace(str(q.get("explanation", "")))

        options_ok = (
            isinstance(options, list)
            and len(options) == expected_n_options
            and all(_clean_whitespace(str(o)) for o in options)
        )
        answer_ok = answer_letter in VALID_ANSWER_LETTERS

        if answer_ok:
            valid_answer_n += 1
        if explanation:
            has_explanation_n += 1
        if q_text and options_ok and answer_ok:
            valid_format_n += 1

    quiz_count = len(questions)
    answer_valid_rate = round((valid_answer_n / quiz_count), 4) if quiz_count else 0.0
    explanation_rate = round((has_explanation_n / quiz_count), 4) if quiz_count else 0.0

    quiz_format_validity = (
        quiz_count == expected_n_questions
        and valid_format_n == expected_n_questions
    )

    return {
        "quiz_count": quiz_count,
        "quiz_format_validity": bool(quiz_format_validity),
        "quiz_answer_valid_rate": answer_valid_rate,
        "quiz_explanation_rate": explanation_rate,
        "quiz_parse_error": loaded["error"],
        "quiz_title": loaded["title"],
        "quiz_questions": questions,
        "quiz_text": _quiz_text_blob(questions),
    }


# ---------------------------------------------------------------------------
# Video
# ---------------------------------------------------------------------------

def _probe_duration(video_path: str) -> Optional[float]:
    try:
        proc = subprocess.run(
            [
                "ffprobe",
                "-v",
                "error",
                "-show_entries",
                "format=duration",
                "-of",
                "default=noprint_wrappers=1:nokey=1",
                video_path,
            ],
            capture_output=True,
            text=True,
            timeout=10,
        )
        if proc.returncode == 0 and proc.stdout.strip():
            return round(float(proc.stdout.strip()), 2)
    except Exception:
        pass

    try:
        from moviepy.editor import VideoFileClip

        clip = VideoFileClip(video_path)
        dur = float(clip.duration)
        clip.close()
        return round(dur, 2)
    except Exception:
        return None


def evaluate_video(video_path: str) -> Dict[str, Any]:
    path = Path(video_path)
    exists = bool(video_path) and path.is_file()
    size_bytes = path.stat().st_size if exists else 0

    return {
        "video_exists": exists,
        "video_size_bytes": size_bytes,
        "video_size_mb": round(size_bytes / (1024 * 1024), 3) if exists else 0.0,
        "video_duration_sec": _probe_duration(str(path)) if exists else None,
        "video_pass": bool(exists and size_bytes > 0),
    }


# ---------------------------------------------------------------------------
# TF-IDF concept extraction
# ---------------------------------------------------------------------------

def _simple_lemmatize(token: str) -> str:
    t = token.lower()
    if t in _IRREGULAR_LEMMA:
        return _IRREGULAR_LEMMA[t]

    if len(t) > 4 and t.endswith("ies"):
        return t[:-3] + "y"
    if len(t) > 5 and t.endswith("ing"):
        base = t[:-3]
        if len(base) >= 2 and base[-1] == base[-2]:
            base = base[:-1]
        return base
    if len(t) > 4 and t.endswith("ed"):
        return t[:-2]
    if len(t) > 4 and t.endswith("es"):
        return t[:-2]
    if len(t) > 3 and t.endswith("s"):
        return t[:-1]
    return t


def _tokenize(text: str) -> List[str]:
    raw = re.findall(r"[a-zA-Z][a-zA-Z0-9_\-]*", (text or "").lower())
    out: List[str] = []
    for tok in raw:
        tok = _simple_lemmatize(tok)
        if not tok or tok in _STOPWORDS or len(tok) <= 1:
            continue
        out.append(tok)
    return out


def _make_ngrams(tokens: Sequence[str], ngram_range: Tuple[int, int] = (1, 2)) -> List[str]:
    lo, hi = ngram_range
    grams: List[str] = []
    for n in range(lo, hi + 1):
        if n <= 0 or len(tokens) < n:
            continue
        for i in range(0, len(tokens) - n + 1):
            grams.append(" ".join(tokens[i : i + n]))
    return grams


def _compute_df(term_docs: Sequence[Sequence[str]]) -> Counter:
    df: Counter = Counter()
    for doc_terms in term_docs:
        df.update(set(doc_terms))
    return df


def _top_tfidf_terms(
    doc_terms: Sequence[str],
    df: Counter,
    n_docs: int,
    top_k: int,
) -> List[str]:
    if not doc_terms:
        return []

    tf = Counter(doc_terms)
    scored: List[Tuple[str, float]] = []

    for term, freq in tf.items():
        term_df = df.get(term, 0)
        idf = math.log((n_docs + 1) / (term_df + 1)) + 1.0
        score = float(freq) * idf
        scored.append((term, score))

    scored.sort(key=lambda x: (-x[1], x[0]))
    return [term for term, _ in scored[:top_k]]


def build_tfidf_keyword_sets(
    texts_by_key: Dict[str, str],
    top_k: int = TFIDF_TOP_K,
    ngram_range: Tuple[int, int] = (1, 2),
) -> Dict[str, set[str]]:
    keys = list(texts_by_key.keys())
    doc_terms: List[List[str]] = []

    for key in keys:
        tokens = _tokenize(texts_by_key.get(key, ""))
        terms = _make_ngrams(tokens, ngram_range=ngram_range)
        doc_terms.append(terms)

    n_docs = max(1, len(doc_terms))
    df = _compute_df(doc_terms)

    out: Dict[str, set[str]] = {}
    for key, terms in zip(keys, doc_terms):
        top_terms = _top_tfidf_terms(terms, df=df, n_docs=n_docs, top_k=top_k)
        out[key] = set(top_terms)
    return out


def compute_align_sc(ks: set[str], kc: set[str]) -> float:
    return round(_jaccard(ks, kc), 4)


def compute_coverage_out_of_scope(kq: set[str], ks: set[str], kc: set[str]) -> Tuple[float, float]:
    km = set(ks) | set(kc)
    if not kq:
        return 0.0, 1.0
    coverage = len(set(kq) & km) / len(set(kq))
    coverage = round(coverage, 4)
    return coverage, round(1.0 - coverage, 4)


def compute_consistency_6pair(keyword_sets: Dict[str, set[str]]) -> float:
    ks = keyword_sets.get("S", set())
    kc = keyword_sets.get("C", set())
    kq = keyword_sets.get("Q", set())
    kv = keyword_sets.get("V", set())

    pairs = [
        _jaccard(ks, kc),
        _jaccard(ks, kq),
        _jaccard(ks, kv),
        _jaccard(kc, kq),
        _jaccard(kc, kv),
        _jaccard(kq, kv),
    ]
    return round(sum(pairs) / len(pairs), 4)


def _keywords_for_segments(segments: Sequence[str], top_k: int = TFIDF_TOP_K) -> List[set[str]]:
    if not segments:
        return []

    term_docs = [_make_ngrams(_tokenize(s), ngram_range=(1, 2)) for s in segments]
    n_docs = len(term_docs)
    df = _compute_df(term_docs)

    out: List[set[str]] = []
    for terms in term_docs:
        out.append(set(_top_tfidf_terms(terms, df=df, n_docs=n_docs, top_k=top_k)))
    return out


def compute_sync_pagewise(
    slide_segments: Sequence[str],
    video_segments: Sequence[str],
    top_k: int = TFIDF_TOP_K,
) -> float:
    n = min(len(slide_segments), len(video_segments))
    if n <= 0:
        return 0.0

    slide_kw = _keywords_for_segments(list(slide_segments)[:n], top_k=top_k)
    video_kw = _keywords_for_segments(list(video_segments)[:n], top_k=top_k)

    sims = [_jaccard(slide_kw[i], video_kw[i]) for i in range(n)]
    return round(sum(sims) / len(sims), 4) if sims else 0.0


def compute_concept_metrics(
    slides_text: str,
    code_text: str,
    quiz_text: str,
    video_text: str,
    slide_segments: Sequence[str],
    video_segments: Sequence[str],
) -> Dict[str, Any]:
    keyword_sets = build_tfidf_keyword_sets(
        {
            "S": slides_text or "",
            "C": code_text or "",
            "Q": quiz_text or "",
            "V": video_text or "",
        },
        top_k=TFIDF_TOP_K,
        ngram_range=(1, 2),
    )

    align_sc = compute_align_sc(keyword_sets["S"], keyword_sets["C"])
    coverage, out_of_scope = compute_coverage_out_of_scope(keyword_sets["Q"], keyword_sets["S"], keyword_sets["C"])
    consistency = compute_consistency_6pair(keyword_sets)
    sync = compute_sync_pagewise(slide_segments, video_segments, top_k=TFIDF_TOP_K)

    return {
        "AlignSC": align_sc,
        "Coverage": coverage,
        "OutOfScope": out_of_scope,
        "Consistency": consistency,
        "Sync": sync,
        "KS": sorted(keyword_sets["S"]),
        "KC": sorted(keyword_sets["C"]),
        "KQ": sorted(keyword_sets["Q"]),
        "KV": sorted(keyword_sets["V"]),
    }
