"""
Deterministic PPT Builder

This module provides reliable, tested functions to build PowerPoint presentations
from a Storyboard JSON. No LLM code generation - pure Python.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt

from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)

# =============================================================================
# SLIDE DELETION (Proper implementation)
# =============================================================================

def delete_all_slides(prs: Presentation) -> None:
    """
    Safely delete all slides from a presentation.
    Must drop relationships AND remove from slide ID list.
    """
    for i in range(len(prs.slides) - 1, -1, -1):
        slide_id = prs.slides._sldIdLst[i]
        rId = slide_id.rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]


# =============================================================================
# VISUAL ASSET GENERATORS
# =============================================================================

def create_code_image(code: str, language: str, output_path: Path) -> Path:
    """
    Render code as an image using PIL.
    Creates a clean, readable code block with syntax-like styling.
    """
    # Configuration
    padding = 20
    line_height = 24
    font_size = 16
    bg_color = (40, 44, 52)  # Dark background (like VS Code dark theme)
    text_color = (171, 178, 191)  # Light gray text
    
    lines = code.split('\n')
    max_line_width = max(len(line) for line in lines) if lines else 20
    
    # Calculate image size
    width = max(400, min(800, max_line_width * 10 + padding * 2))
    height = len(lines) * line_height + padding * 2
    
    # Create image
    img = Image.new('RGB', (width, height), color=bg_color)
    draw = ImageDraw.Draw(img)
    
    # Try to use a monospace font, fallback to default
    try:
        font = ImageFont.truetype("consola.ttf", font_size)
    except:
        try:
            font = ImageFont.truetype("DejaVuSansMono.ttf", font_size)
        except:
            font = ImageFont.load_default()
    
    # Draw each line
    y = padding
    for line in lines:
        draw.text((padding, y), line, fill=text_color, font=font)
        y += line_height
    
    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(str(output_path))
    logger.info(f"Created code image: {output_path}")
    return output_path


def create_chart_image(chart_data: Dict[str, Any], output_path: Path) -> Path:
    """
    Generate a chart image using matplotlib.
    Supports: bar, line, pie charts.
    """
    chart_type = chart_data.get("chart_type", "bar")
    title = chart_data.get("title", "Chart")
    data = chart_data.get("data", {})
    labels = data.get("labels", [])
    values = data.get("values", [])
    
    # Create figure
    fig, ax = plt.subplots(figsize=(6, 4))
    
    # Professional colors
    colors = ['#4285F4', '#EA4335', '#FBBC05', '#34A853', '#9C27B0', '#FF5722']
    
    if chart_type == "bar":
        bars = ax.bar(labels, values, color=colors[:len(labels)])
        ax.set_ylabel('Value')
    elif chart_type == "line":
        ax.plot(labels, values, marker='o', linewidth=2, color=colors[0])
        ax.set_ylabel('Value')
    elif chart_type == "pie":
        ax.pie(values, labels=labels, autopct='%1.1f%%', colors=colors[:len(labels)])
    else:
        # Default to bar
        ax.bar(labels, values, color=colors[:len(labels)])
    
    ax.set_title(title, fontsize=14, fontweight='bold')
    
    # Clean styling
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    
    plt.tight_layout()
    
    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(str(output_path), dpi=150, bbox_inches='tight', 
                facecolor='white', edgecolor='none')
    plt.close(fig)
    
    logger.info(f"Created chart image: {output_path}")
    return output_path


def create_formula_image(latex: str, output_path: Path) -> Path:
    """
    Render a LaTeX formula as an image using matplotlib.
    """
    fig = plt.figure(figsize=(6, 1.5))
    
    # Clean the latex string
    latex_clean = latex.replace('\\\\', '\\')
    
    # Render the formula
    fig.text(0.5, 0.5, f"${latex_clean}$", 
             fontsize=24, ha='center', va='center',
             transform=fig.transFigure)
    
    plt.axis('off')
    
    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(str(output_path), dpi=200, bbox_inches='tight',
                facecolor='white', edgecolor='none', pad_inches=0.2)
    plt.close(fig)
    
    logger.info(f"Created formula image: {output_path}")
    return output_path


# =============================================================================
# MAIN BUILDER FUNCTION
# =============================================================================

def build_presentation(
    storyboard: Dict[str, Any],
    template_path: Path,
    output_path: Path,
    assets_dir: Path
) -> Path:
    """
    Build a PowerPoint presentation from a Storyboard JSON.
    
    Args:
        storyboard: The storyboard dict with 'slides' list
        template_path: Path to the master template .pptx
        output_path: Where to save the generated .pptx
        assets_dir: Directory to save generated images
        
    Returns:
        Path to the generated .pptx file
    """
    logger.info(f"Building presentation from storyboard...")
    
    # Ensure assets directory exists
    assets_dir.mkdir(parents=True, exist_ok=True)
    
    # Load template
    prs = Presentation(str(template_path))
    
    # Delete existing slides
    delete_all_slides(prs)
    
    # Get available layouts
    # Layout 0: Title Slide
    # Layout 5: 1-Column Content
    # Layout 6: 2-Column Content
    
    slides_data = storyboard.get("slides", [])
    
    for idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx", 5)
        
        # Clamp layout_idx to valid range
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = min(5, len(prs.slide_layouts) - 1)
        
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        
        # --- Handle Title Slide (Layout 0) ---
        if layout_idx == 0:
            _populate_title_slide(slide, slide_data)
        
        # --- Handle 1-Column Slide (Layout 5) ---
        elif layout_idx == 5:
            _populate_content_slide(slide, slide_data)
        
        # --- Handle 2-Column Slide (Layout 6) ---
        elif layout_idx == 6:
            _populate_two_column_slide(slide, slide_data, assets_dir, idx)
        
        # --- Fallback for other layouts ---
        else:
            _populate_content_slide(slide, slide_data)
    
    # Save presentation
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    logger.info(f"Presentation saved to: {output_path}")
    return output_path


def _populate_title_slide(slide, slide_data: Dict[str, Any]) -> None:
    """Populate a title slide (Layout 0)."""
    title_text = slide_data.get("title", "Untitled")
    subtitle_text = slide_data.get("subtitle", "")
    
    # Find title placeholder (usually index 0)
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.placeholder_format.idx == 0:
                # Title
                shape.text_frame.paragraphs[0].text = title_text
                _style_paragraph(shape.text_frame.paragraphs[0])
            elif shape.placeholder_format.idx == 1:
                # Subtitle
                shape.text_frame.paragraphs[0].text = subtitle_text
                _style_paragraph(shape.text_frame.paragraphs[0])


def _populate_content_slide(slide, slide_data: Dict[str, Any]) -> None:
    """Populate a 1-column content slide (Layout 5)."""
    title_text = slide_data.get("title", "")
    content_text = slide_data.get("content", "")
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        ph_idx = shape.placeholder_format.idx if hasattr(shape, 'placeholder_format') and shape.placeholder_format else None
        
        if ph_idx == 0:
            # Title
            shape.text_frame.paragraphs[0].text = title_text
            _style_paragraph(shape.text_frame.paragraphs[0])
        elif ph_idx in [1, 10, 11, 12, 13, 14]:
            # Body - try common placeholder indices
            shape.text_frame.paragraphs[0].text = content_text
            _style_paragraph(shape.text_frame.paragraphs[0])


def _populate_two_column_slide(
    slide, 
    slide_data: Dict[str, Any], 
    assets_dir: Path,
    slide_idx: int
) -> None:
    """Populate a 2-column slide (Layout 6) with text and visual."""
    title_text = slide_data.get("title", "")
    content_text = slide_data.get("content", "")
    visual_assets = slide_data.get("visual_assets", [])
    
    # Track which placeholders we've found
    title_shape = None
    left_shape = None
    right_shape = None
    
    for shape in slide.shapes:
        if not hasattr(shape, 'placeholder_format') or not shape.placeholder_format:
            continue
            
        ph_idx = shape.placeholder_format.idx
        
        if ph_idx == 0:
            title_shape = shape
        elif ph_idx in [1, 12]:
            # Left column (text)
            left_shape = shape
        elif ph_idx in [2, 13]:
            # Right column (visual)
            right_shape = shape
    
    # Populate title
    if title_shape and title_shape.has_text_frame:
        title_shape.text_frame.paragraphs[0].text = title_text
        _style_paragraph(title_shape.text_frame.paragraphs[0])
    
    # Populate left column (text)
    if left_shape and left_shape.has_text_frame:
        left_shape.text_frame.paragraphs[0].text = content_text
        _style_paragraph(left_shape.text_frame.paragraphs[0])
    
    # Populate right column (visual)
    if visual_assets and right_shape:
        visual = visual_assets[0]  # Use first visual asset
        visual_type = visual.get("type", "")
        
        # Generate the image
        img_path = None
        
        if visual_type == "code_snippet":
            code = visual.get("code", "# No code")
            lang = visual.get("language", "python")
            img_path = assets_dir / f"code_{slide_idx}.png"
            create_code_image(code, lang, img_path)
            
        elif visual_type == "chart_data":
            img_path = assets_dir / f"chart_{slide_idx}.png"
            create_chart_image(visual, img_path)
            
        elif visual_type == "formula_latex":
            latex = visual.get("content", "E = mc^2")
            img_path = assets_dir / f"formula_{slide_idx}.png"
            create_formula_image(latex, img_path)
        
        # Insert the image
        if img_path and img_path.exists():
            # Get the placeholder position and size
            left = right_shape.left
            top = right_shape.top
            width = right_shape.width
            height = right_shape.height
            
            # Remove the placeholder shape and add picture
            sp = right_shape._element
            sp.getparent().remove(sp)
            
            # Add the picture
            slide.shapes.add_picture(str(img_path), left, top, width, height)


def _style_paragraph(paragraph) -> None:
    """Apply consistent styling to a paragraph."""
    for run in paragraph.runs:
        run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

