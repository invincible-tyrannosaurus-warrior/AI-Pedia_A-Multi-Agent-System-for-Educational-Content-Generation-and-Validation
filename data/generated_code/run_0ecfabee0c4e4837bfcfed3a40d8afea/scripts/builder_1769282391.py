import json
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt

# Helper function to create LaTeX formula images
def create_latex_image(formula, output_path):
    fig = plt.figure(figsize=(3, 1))
    text = fig.text(0.5, 0.5, f"${formula}$", fontsize=20, ha='center', va='center')
    plt.axis('off')
    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close()

# Helper function to create code snippet images
def create_code_snippet_image(code, output_path):
    # Create an image with white background
    img = Image.new('RGB', (600, 400), color='white')
    draw = ImageDraw.Draw(img)
    
    # Load a font (fallback to default if not available)
    try:
        font = ImageFont.truetype("arial.ttf", 16)
    except:
        font = ImageFont.load_default()
    
    # Split code into lines and draw each line
    lines = code.split('\n')
    y_position = 10
    for line in lines:
        draw.text((10, y_position), line, fill='black', font=font)
        y_position += 25
    
    img.save(output_path)

# Function to create bar chart
def create_bar_chart(data, output_path):
    labels = data.get('labels', [])
    values = data.get('values', [])
    
    plt.figure(figsize=(8, 6))
    plt.style.use('ggplot')
    bars = plt.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
    
    # Add value labels on top of bars
    for bar, value in zip(bars, values):
        plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.01,
                str(value), ha='center', va='bottom')
    
    plt.title(data.get('title', 'Chart'))
    plt.xlabel(data.get('x_label', 'X Axis'))
    plt.ylabel(data.get('y_label', 'Y Axis'))
    plt.xticks(rotation=45)
    plt.tight_layout()
    
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

# Function to create pie chart
def create_pie_chart(data, output_path):
    labels = data.get('labels', [])
    values = data.get('values', [])
    
    plt.figure(figsize=(8, 6))
    plt.style.use('ggplot')
    colors = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd']
    wedges, texts, autotexts = plt.pie(values, labels=labels, autopct='%1.1f%%', startangle=90, colors=colors)
    
    plt.title(data.get('title', 'Pie Chart'))
    plt.axis('equal')
    plt.tight_layout()
    
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

# Function to create line chart
def create_line_chart(data, output_path):
    x_values = data.get('x_values', [])
    y_values = data.get('y_values', [])
    
    plt.figure(figsize=(8, 6))
    plt.style.use('ggplot')
    plt.plot(x_values, y_values, marker='o', color='#1f77b4')
    
    plt.title(data.get('title', 'Line Chart'))
    plt.xlabel(data.get('x_label', 'X Axis'))
    plt.ylabel(data.get('y_label', 'Y Axis'))
    plt.grid(True)
    plt.tight_layout()
    
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

# Load storyboard
with open(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_0ecfabee0c4e4837bfcfed3a40d8afea/storyboard.json", "r", encoding="utf-8") as f:
    storyboard = json.load(f)

slides_data = storyboard.get("slides", [])

# Load presentation template
template_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx"
prs = Presentation(template_path)

# Remove existing slides using the correct method
# Remove slides by iterating backwards through the list
for i in range(len(prs.slides) - 1, -1, -1):
    prs.slides._sldIdLst.remove(prs.slides[i]._element)

# Process each slide
for slide_idx, slide_data in enumerate(slides_data):
    slide_type = slide_data.get("type")
    
    # Determine layout based on slide type
    if slide_type == "title":
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        title = slide.placeholders[0]
        subtitle = slide.placeholders[1]
        
        title.text = slide_data.get("title", "")
        subtitle.text = slide_data.get("subtitle", "")
        
    elif slide_type == "content_single_column":
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Content Slide (1 Column)
        title = slide.placeholders[0]
        body = slide.placeholders[11]
        
        title.text = slide_data.get("title", "")
        
        # Set text color to black
        for paragraph in body.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)
        
        body.text = slide_data.get("content", "")
        
    elif slide_type == "content_double_column":
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Content Slide (2 Columns)
        title = slide.placeholders[0]
        left_column = slide.placeholders[12]
        right_column = slide.placeholders[13]
        
        title.text = slide_data.get("title", "")
        
        # Set text color to black for left column
        for paragraph in left_column.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)
                
        left_column.text = slide_data.get("left_content", "")
        
        # Handle visual assets in right column
        visual_assets = slide_data.get("visual_assets", [])
        if visual_assets:
            asset = visual_assets[0]  # Take first visual asset
            
            # Check if it's a path to an existing image
            if isinstance(asset, str) and os.path.exists(asset):
                try:
                    # Get placeholder dimensions
                    ph = right_column
                    slide.shapes.add_picture(asset, ph.left, ph.top, ph.width, ph.height)
                except Exception as e:
                    print(f"Error inserting image {asset}: {e}")
            
            # Handle code snippet
            elif isinstance(asset, dict) and asset.get("type") == "code_snippet":
                code = asset.get("content", "")
                img_path = f"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/assets/code_{slide_idx}.png"
                try:
                    create_code_snippet_image(code, img_path)
                    ph = right_column
                    slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                except Exception as e:
                    print(f"Error creating code snippet image: {e}")
            
            # Handle chart data
            elif isinstance(asset, dict) and asset.get("type") == "chart_data":
                chart_data = asset.get("data", {})
                chart_type = asset.get("chart_type", "bar")
                img_path = f"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/assets/chart_{slide_idx}.png"
                
                try:
                    if chart_type == "bar":
                        create_bar_chart(chart_data, img_path)
                    elif chart_type == "pie":
                        create_pie_chart(chart_data, img_path)
                    elif chart_type == "line":
                        create_line_chart(chart_data, img_path)
                    else:
                        raise ValueError(f"Unsupported chart type: {chart_type}")
                    
                    ph = right_column
                    slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                except Exception as e:
                    print(f"Error creating chart: {e}")
                    # Create error image
                    error_img_path = f"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/assets/chart_error_{slide_idx}.png"
                    create_code_snippet_image("Chart Error", error_img_path)
                    ph = right_column
                    slide.shapes.add_picture(error_img_path, ph.left, ph.top, ph.width, ph.height)
            
            # Handle formula
            elif isinstance(asset, dict) and asset.get("type") == "formula_latex":
                formula = asset.get("content", "")
                img_path = f"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/assets/formula_{slide_idx}.png"
                try:
                    create_latex_image(formula, img_path)
                    ph = right_column
                    slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                except Exception as e:
                    print(f"Error creating formula image: {e}")

# Save the presentation
output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_0ecfabee0c4e4837bfcfed3a40d8afea/lesson_1769282391.pptx"
prs.save(output_path)
