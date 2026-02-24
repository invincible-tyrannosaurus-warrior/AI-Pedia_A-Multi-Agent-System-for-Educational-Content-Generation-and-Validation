import json
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt

def create_code_image(code_text, output_path):
    # Create an image with code snippet
    img = Image.new('RGB', (800, 600), color='white')
    draw = ImageDraw.Draw(img)
    
    # Load a font (fallback to default if not available)
    try:
        font = ImageFont.truetype("cour.ttf", 16)
    except:
        font = ImageFont.load_default()
    
    # Draw text
    draw.text((10, 10), code_text, fill='black', font=font)
    
    # Save image
    img.save(output_path)

def create_latex_image(formula, output_path):
    fig = plt.figure(figsize=(3, 1))
    text = fig.text(0.5, 0.5, f"${formula}$", fontsize=20, ha='center', va='center')
    plt.axis('off')
    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close()

def create_chart_image(chart_data, output_path):
    try:
        labels = chart_data.get('labels', [])
        values = chart_data.get('values', [])
        chart_type = chart_data.get('type', 'bar')
        
        # Set up the figure
        fig, ax = plt.subplots(figsize=(8, 6))
        fig.patch.set_alpha(0)
        ax.patch.set_alpha(0)
        
        # Create chart based on type
        if chart_type == 'bar':
            bars = ax.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
        elif chart_type == 'line':
            ax.plot(labels, values, marker='o', color='#1f77b4')
        elif chart_type == 'pie':
            ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=90)
        else:
            bars = ax.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
            
        # Customize chart
        ax.set_facecolor('white')
        fig.patch.set_facecolor('white')
        
        # Save chart
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
    except Exception as e:
        # If chart creation fails, create error image
        img = Image.new('RGB', (400, 200), color='white')
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("cour.ttf", 16)
        except:
            font = ImageFont.load_default()
        draw.text((10, 10), f"Chart Error: {str(e)}", fill='red', font=font)
        img.save(output_path)

def main():
    # Define paths
    storyboard_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_0ecfabee0c4e4837bfcfed3a40d8afea/storyboard.json"
    assets_dir = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/assets"
    template_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx"
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_0ecfabee0c4e4837bfcfed3a40d8afea/lesson_1769282480.pptx"
    
    # Load storyboard
    with open(storyboard_path, "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(template_path)
    
    # Remove existing slides using the correct method
    # Create a new presentation with only the title slide
    new_prs = Presentation(template_path)
    prs = new_prs
    
    # Process each slide
    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        visual_assets = slide_data.get("visual_assets", [])
        
        # Determine layout based on slide type
        if slide_type == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
            title_placeholder = slide.placeholders[0]
            subtitle_placeholder = slide.placeholders[1]
            
            title_placeholder.text = title
            subtitle_placeholder.text = content
            
        elif slide_type == "content_single_column":
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Content Slide (1 Column)
            title_placeholder = slide.placeholders[0]
            body_placeholder = slide.placeholders[11]
            
            title_placeholder.text = title
            body_placeholder.text = content
            
        elif slide_type == "content_double_column":
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Content Slide (2 Columns)
            title_placeholder = slide.placeholders[0]
            left_placeholder = slide.placeholders[12]
            right_placeholder = slide.placeholders[13]
            
            title_placeholder.text = title
            
            # Add text to left column
            left_paragraph = left_placeholder.text_frame.paragraphs[0]
            left_paragraph.text = content
            left_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # Handle visual assets in right column
            if visual_assets:
                asset = visual_assets[0]  # Take first visual asset
                
                if isinstance(asset, str):
                    # It's a path to an image
                    if os.path.exists(asset):
                        # Get the dimensions of the placeholder
                        ph = right_placeholder
                        slide.shapes.add_picture(asset, ph.left, ph.top, ph.width, ph.height)
                elif isinstance(asset, dict):
                    # It's a structured asset
                    asset_type = asset.get("type", "")
                    
                    if asset_type == "code_snippet":
                        code_content = asset.get("content", "")
                        img_path = os.path.join(assets_dir, f"code_{len(prs.slides)}.png")
                        create_code_image(code_content, img_path)
                        ph = right_placeholder
                        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                        
                    elif asset_type == "chart_data":
                        chart_data = asset.get("data", {})
                        img_path = os.path.join(assets_dir, f"chart_{len(prs.slides)}.png")
                        create_chart_image(chart_data, img_path)
                        ph = right_placeholder
                        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                        
                    elif asset_type == "formula_latex":
                        formula = asset.get("content", "")
                        img_path = os.path.join(assets_dir, f"formula_{len(prs.slides)}.png")
                        create_latex_image(formula, img_path)
                        ph = right_placeholder
                        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
    
    # Save presentation
    prs.save(output_path)

if __name__ == "__main__":
    main()
