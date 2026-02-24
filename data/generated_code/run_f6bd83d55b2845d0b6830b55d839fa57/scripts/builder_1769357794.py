import json
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt

def delete_all_slides(prs):
    # Must iterate in reverse to avoid index shifting issues
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

def create_code_snippet_image(code, output_path):
    # Create an image with code snippet
    img = Image.new('RGB', (600, 400), color='white')
    draw = ImageDraw.Draw(img)
    
    # Try to use a monospace font if available
    try:
        font = ImageFont.truetype("cour.ttf", 16)
    except:
        font = ImageFont.load_default()
    
    # Draw code text
    draw.text((10, 10), code, fill='black', font=font)
    
    # Save image
    img.save(output_path)

def create_latex_image(formula, output_path):
    fig = plt.figure(figsize=(3, 1))
    text = fig.text(0.5, 0.5, f"${formula}$", fontsize=20, ha='center', va='center')
    plt.axis('off')
    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close()

def create_bar_chart(data, output_path):
    labels = data.get('labels', [])
    values = data.get('values', [])
    
    # Create figure with white background
    fig, ax = plt.subplots(figsize=(8, 6))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    
    # Create bar chart
    bars = ax.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
    
    # Customize chart
    ax.set_xlabel('Categories')
    ax.set_ylabel('Values')
    ax.set_title('Chart')
    
    # Rotate x-axis labels if needed
    plt.xticks(rotation=45, ha='right')
    
    # Adjust layout to prevent label cutoff
    plt.tight_layout()
    
    # Save chart
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

def create_line_chart(data, output_path):
    labels = data.get('labels', [])
    values = data.get('values', [])
    
    # Create figure with white background
    fig, ax = plt.subplots(figsize=(8, 6))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    
    # Create line chart
    ax.plot(labels, values, marker='o', color='#1f77b4')
    
    # Customize chart
    ax.set_xlabel('Categories')
    ax.set_ylabel('Values')
    ax.set_title('Line Chart')
    
    # Rotate x-axis labels if needed
    plt.xticks(rotation=45, ha='right')
    
    # Adjust layout to prevent label cutoff
    plt.tight_layout()
    
    # Save chart
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

def create_pie_chart(data, output_path):
    labels = data.get('labels', [])
    values = data.get('values', [])
    
    # Create figure with white background
    fig, ax = plt.subplots(figsize=(8, 6))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    
    # Create pie chart
    wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=90)
    
    # Customize chart
    ax.set_title('Pie Chart')
    
    # Adjust layout to prevent label cutoff
    plt.tight_layout()
    
    # Save chart
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

def process_visual_asset(asset, asset_idx, assets_dir):
    if isinstance(asset, str):
        # Case A: String path
        if os.path.exists(asset):
            return asset
        else:
            return None
    elif asset.get('type') == 'code_snippet':
        # Case B: Code snippet
        code = asset.get('content', '')
        output_path = os.path.join(assets_dir, f"code_{asset_idx}.png")
        create_code_snippet_image(code, output_path)
        return output_path
    elif asset.get('type') == 'chart_data':
        # Case C: Chart
        chart_type = asset.get('chart_type', 'bar')
        data = asset.get('data', {})
        output_path = os.path.join(assets_dir, f"chart_{asset_idx}.png")
        
        try:
            if chart_type == 'bar':
                create_bar_chart(data, output_path)
            elif chart_type == 'line':
                create_line_chart(data, output_path)
            elif chart_type == 'pie':
                create_pie_chart(data, output_path)
            else:
                # Default to bar chart
                create_bar_chart(data, output_path)
            return output_path
        except Exception as e:
            # If chart creation fails, create error image
            print(f"Error creating chart: {e}")
            error_img_path = os.path.join(assets_dir, f"chart_{asset_idx}_error.png")
            create_code_snippet_image("Chart Error", error_img_path)
            return error_img_path
    elif asset.get('type') == 'formula_latex':
        # Case D: Formula
        formula = asset.get('content', '')
        output_path = os.path.join(assets_dir, f"formula_{asset_idx}.png")
        create_latex_image(formula, output_path)
        return output_path
    else:
        return None

# Load storyboard
with open(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_f6bd83d55b2845d0b6830b55d839fa57/storyboard.json", "r", encoding="utf-8") as f:
    storyboard = json.load(f)

slides_data = storyboard.get("slides", [])

# Load presentation template
template_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx"
prs = Presentation(template_path)

# Delete all existing slides
delete_all_slides(prs)

# Define assets directory
assets_dir = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_f6bd83d55b2845d0b6830b55d839fa57/assets"

# Process each slide
for slide_data in slides_data:
    slide_type = slide_data.get("type", "content")
    
    if slide_type == "title":
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
        title = slide.placeholders[0]
        subtitle = slide.placeholders[1]
        
        title.text = slide_data.get("title", "")
        subtitle.text = slide_data.get("subtitle", "")
        
    elif slide_type == "content":
        content_layout = slide_data.get("layout", "single_column")
        
        if content_layout == "single_column":
            # Content slide (1 column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Content Slide (1 Column)
            title = slide.placeholders[0]
            body = slide.placeholders[11]
            
            title.text = slide_data.get("title", "")
            
            # Set body text with black color
            body_text = body.text_frame
            body_text.clear()
            p = body_text.add_paragraph()
            p.text = slide_data.get("content", "")
            p.font.color.rgb = RGBColor(0, 0, 0)
            
        elif content_layout == "two_columns":
            # Content slide (2 columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Content Slide (2 Columns)
            title = slide.placeholders[0]
            left_column = slide.placeholders[12]
            right_column = slide.placeholders[13]
            
            title.text = slide_data.get("title", "")
            
            # Set left column text with black color
            left_text = left_column.text_frame
            left_text.clear()
            p = left_text.add_paragraph()
            p.text = slide_data.get("content", "")
            p.font.color.rgb = RGBColor(0, 0, 0)
            
            # Process visual assets for right column
            visual_assets = slide_data.get("visual_assets", [])
            if visual_assets:
                asset = visual_assets[0]  # Take first asset
                asset_path = process_visual_asset(asset, 0, assets_dir)
                
                if asset_path:
                    # Add picture to right column
                    try:
                        # Get dimensions of placeholder
                        ph = right_column
                        slide.shapes.add_picture(
                            asset_path,
                            ph.left,
                            ph.top,
                            ph.width,
                            ph.height
                        )
                    except Exception as e:
                        print(f"Error inserting image: {e}")

# Save presentation
output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_f6bd83d55b2845d0b6830b55d839fa57/lesson_1769357794.pptx"
prs.save(output_path)
