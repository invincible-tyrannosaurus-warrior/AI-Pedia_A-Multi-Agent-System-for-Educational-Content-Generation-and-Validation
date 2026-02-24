from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import json
import os

def build_presentation():
    # Load storyboard
    with open(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_82e0d4a3990b4419894ac251f01bf602/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx")
        
        # Add new slide based on layout
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            try:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
            except KeyError:
                pass
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            
            # Set body text with black color for all paragraphs
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[11]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            
            # Set left column content with black color for all paragraphs
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[12]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Handle visual assets in right column
            visual_assets = slide_data.get("visual_assets", [])
            if visual_assets:
                asset = visual_assets[0] if isinstance(visual_assets, list) else visual_assets
                
                # Case B: Code snippet dictionary
                if isinstance(asset, dict) and asset.get("type") == "code_snippet":
                    code = asset.get("code", "")
                    
                    # Create image with code
                    img = Image.new('RGB', (800, 600), 'white')
                    draw = ImageDraw.Draw(img)
                    
                    # Try to use a monospace font if available
                    try:
                        font = ImageFont.truetype("cour.ttf", 16)
                    except:
                        font = ImageFont.load_default()
                    
                    # Split code into lines and draw
                    lines = code.split('\n')
                    y_text = 10
                    for line in lines:
                        draw.text((10, y_text), line, fill='black', font=font)
                        y_text += 20
                    
                    # Save image
                    img_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_82e0d4a3990b4419894ac251f01bf602/assets/generated_code_{}.png".format(slide_idx)
                    img.save(img_path)
                    
                    # Insert image into placeholder
                    try:
                        ph = slide.placeholders[13]
                        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                    except KeyError:
                        pass
                        
                # Case A: File path string
                elif isinstance(asset, str):
                    if os.path.exists(asset):
                        try:
                            ph = slide.placeholders[13]
                            slide.shapes.add_picture(asset, ph.left, ph.top, ph.width, ph.height)
                        except KeyError:
                            pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_82e0d4a3990b4419894ac251f01bf602/lesson_1769088879.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
