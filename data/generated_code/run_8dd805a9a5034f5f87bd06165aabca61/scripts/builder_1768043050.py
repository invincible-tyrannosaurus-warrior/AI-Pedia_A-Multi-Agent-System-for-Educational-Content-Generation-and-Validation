from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import json
import os

def build_presentation():
    # Load storyboard
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_8dd805a9a5034f5f87bd06165aabca61/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx")
        
        # Add new slide based on layout
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide_data.get("title", "")
            subtitle = slide_data.get("subtitle", "")
            
            try:
                slide.placeholders[0].text = title
            except KeyError:
                pass
            
            try:
                slide.placeholders[1].text = subtitle
            except KeyError:
                pass
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            
            try:
                slide.placeholders[0].text = title
            except KeyError:
                pass
            
            try:
                shape = slide.placeholders[11]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            visual_assets = slide_data.get("visual_assets", [])
            
            # Set title
            try:
                slide.placeholders[0].text = title
            except KeyError:
                pass
            
            # Set left column content
            try:
                shape = slide.placeholders[12]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Handle right column visual assets
            try:
                placeholder = slide.placeholders[13]
                
                # If visual_assets is a list, get first item
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset = visual_assets[0]
                else:
                    asset = visual_assets
                    
                # Case A: String path to image file
                if isinstance(asset, str) and os.path.isfile(asset):
                    placeholder.insert_picture(asset)
                    
                # Case B: Dictionary representing code snippet
                elif isinstance(asset, dict) and asset.get('type') == 'code_snippet':
                    code = asset.get('code', '')
                    
                    # Generate image with PIL
                    img = Image.new('RGB', (800, 600), 'white')
                    draw = ImageDraw.Draw(img)
                    font_size = 16
                    # Simple approach - draw text directly without font loading
                    draw.text((10, 10), code, fill='black')
                    
                    # Save to assets directory
                    img_path = f"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_8dd805a9a5034f5f87bd06165aabca61/assets/generated_code_{slide_idx}.png"
                    img.save(img_path)
                    
                    # Insert into slide
                    placeholder.insert_picture(img_path)
                    
            except KeyError:
                pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768043050.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
