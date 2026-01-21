from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_8dd805a9a5034f5f87bd06165aabca61/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx", 5)
        
        # Add new slide based on layout
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Set title if available
        title = slide_data.get("title", "")
        try:
            slide.placeholders[0].text = title
        except KeyError:
            pass
        
        # Handle different layouts
        if layout_idx == 5:  # Content Slide (1 Column)
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[11]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            # Left column content
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[12]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Right column visual assets
            visual_assets = slide_data.get("visual_assets", None)
            if visual_assets:
                try:
                    placeholder = slide.placeholders[13]
                    
                    # Handle case when visual_assets is a list
                    if isinstance(visual_assets, list):
                        visual_asset = visual_assets[0]
                    else:
                        visual_asset = visual_assets
                    
                    # Case A: String path to image
                    if isinstance(visual_asset, str) and os.path.isfile(visual_asset):
                        placeholder.insert_picture(visual_asset)
                    
                    # Case B: Dictionary representing code snippet
                    elif isinstance(visual_asset, dict) and visual_asset.get("type") == "code_snippet":
                        code = visual_asset.get("code", "")
                        
                        # Generate image with code snippet
                        img_path = f"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_8dd805a9a5034f5f87bd06165aabca61/assets/generated_code_{slide_idx}.png"
                        
                        # Create image with PIL
                        img = Image.new('RGB', (800, 600), 'white')
                        draw = ImageDraw.Draw(img)
                        
                        # Try to use a monospace font if available
                        try:
                            font = ImageFont.truetype("cour.ttf", 16)
                        except:
                            font = ImageFont.load_default()
                        
                        # Draw code text
                        draw.text((10, 10), code, fill='black', font=font)
                        
                        # Save image
                        img.save(img_path)
                        
                        # Insert into slide
                        placeholder.insert_picture(img_path)
                        
                except KeyError:
                    pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768043084.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
