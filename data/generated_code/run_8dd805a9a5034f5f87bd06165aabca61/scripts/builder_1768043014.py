from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import json
import os

def build_presentation():
    # Load storyboard
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_8dd805a9a5034f5f87bd06165aabca61/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])

    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")

    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]

    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx")
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        visual_assets = slide_data.get("visual_assets", [])

        # Create new slide based on layout
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            try:
                slide.placeholders[0].text = title
            except KeyError:
                pass
            try:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
            except KeyError:
                pass
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            try:
                slide.placeholders[0].text = title
            except KeyError:
                pass
            try:
                slide.placeholders[11].text = content
                # Set text color to black for all paragraphs
                for p in slide.placeholders[11].text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            try:
                slide.placeholders[0].text = title
            except KeyError:
                pass
            
            # Left column - content
            try:
                slide.placeholders[12].text = content
                # Set text color to black for all paragraphs
                for p in slide.placeholders[12].text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
            # Right column - visual assets
            try:
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset = visual_assets[0]
                    
                    # Handle string paths
                    if isinstance(asset, str):
                        if os.path.exists(asset):
                            slide.placeholders[13].insert_picture(asset)
                            
                    # Handle dictionary (code snippet)
                    elif isinstance(asset, dict) and asset.get('type') == 'code_snippet':
                        code = asset.get('code', '')
                        
                        # Generate image with PIL
                        img = Image.new('RGB', (800, 600), 'white')
                        draw = ImageDraw.Draw(img)
                        font_size = 16
                        # Simple monospace font approach
                        # In practice, you might want to use a proper monospace font
                        try:
                            from PIL import ImageFont
                            font = ImageFont.truetype("cour.ttf", font_size)
                        except:
                            font = ImageFont.load_default()
                            
                        # Draw text
                        y_text = 10
                        for line in code.split('\n'):
                            draw.text((10, y_text), line, fill='black', font=font)
                            y_text += font_size + 2
                        
                        # Save image
                        img_path = f"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_8dd805a9a5034f5f87bd06165aabca61/assets/generated_code_{slide_idx}.png"
                        img.save(img_path)
                        
                        # Insert picture into slide
                        slide.placeholders[13].insert_picture(img_path)
                        
            except KeyError:
                pass

    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768043014.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
