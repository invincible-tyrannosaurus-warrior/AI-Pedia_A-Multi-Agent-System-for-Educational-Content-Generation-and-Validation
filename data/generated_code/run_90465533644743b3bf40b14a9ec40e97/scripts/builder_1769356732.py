import json
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt

def create_latex_image(formula, output_path):
    """Create a LaTeX formula image using matplotlib."""
    fig = plt.figure(figsize=(3, 1))
    text = fig.text(0.5, 0.5, f"${formula}$", fontsize=20, ha='center', va='center')
    plt.axis('off')
    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close()

def create_code_snippet_image(code, output_path):
    """Create an image of a code snippet using PIL."""
    # Create a white background image
    img = Image.new('RGB', (600, 400), color='white')
    draw = ImageDraw.Draw(img)
    
    # Load a font (fallback to default if not available)
    try:
        font = ImageFont.truetype("arial.ttf", 16)
    except:
        font = ImageFont.load_default()
    
    # Draw the code text
    draw.text((10, 10), code, fill='black', font=font)
    
    # Save the image
    img.save(output_path)

def create_bar_chart(data, output_path):
    """Create a bar chart from data."""
    labels = data.get('labels', [])
    values = data.get('values', [])
    
    # Set up the figure with a clean style
    plt.style.use('ggplot')
    fig, ax = plt.subplots(figsize=(8, 6))
    
    # Create the bar chart
    bars = ax.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
    
    # Customize the chart
    ax.set_xlabel('Categories')
    ax.set_ylabel('Values')
    ax.set_title('Data Visualization')
    
    # Add value labels on bars
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.01,
                str(value), ha='center', va='bottom')
    
    # Adjust layout and save
    plt.tight_layout()
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

def delete_all_slides(prs):
    """Delete all slides from presentation."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for s in slides:
        try:
            xml_slides.remove(s)
        except ValueError:
            pass

def main():
    # Define paths
    storyboard_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_90465533644743b3bf40b14a9ec40e97/storyboard.json"
    assets_dir = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_90465533644743b3bf40b14a9ec40e97/assets"
    template_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx"
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_90465533644743b3bf40b14a9ec40e97/lesson_1769356732.pptx"
    
    # Load storyboard
    with open(storyboard_path, "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template and delete existing slides
    prs = Presentation(template_path)
    delete_all_slides(prs)
    
    # Process each slide
    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type", "content")
        
        if slide_type == "title":
            # Title slide (Layout Index: 0)
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.placeholders[0]
            subtitle = slide.placeholders[1]
            
            title.text = slide_data.get("title", "")
            subtitle.text = slide_data.get("subtitle", "")
            
        elif slide_type == "content":
            # Content slide (1 column) (Layout Index: 5)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.placeholders[0]
            body = slide.placeholders[11]
            
            title.text = slide_data.get("title", "")
            
            # Set text color to black
            for paragraph in body.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add content
            body.text = slide_data.get("content", "")
            
        elif slide_type == "content_2col":
            # Content slide (2 columns) (Layout Index: 6)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.placeholders[0]
            left_col = slide.placeholders[12]
            right_col = slide.placeholders[13]
            
            title.text = slide_data.get("title", "")
            
            # Set text color to black for left column
            for paragraph in left_col.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add content to left column
            left_col.text = slide_data.get("left_column", "")
            
            # Handle visual assets in right column
            visual_assets = slide_data.get("visual_assets", [])
            if visual_assets:
                asset = visual_assets[0]  # Take the first asset
                
                if isinstance(asset, str):
                    # Case A: String path to existing image
                    if os.path.exists(asset):
                        right_col.clear()
                        pic = right_col.insert_picture(asset)
                        
                elif isinstance(asset, dict):
                    # Case B: Code snippet
                    if asset.get("type") == "code_snippet":
                        code = asset.get("content", "")
                        img_path = os.path.join(assets_dir, f"code_{i}.png")
                        create_code_snippet_image(code, img_path)
                        right_col.clear()
                        pic = right_col.insert_picture(img_path)
                        
                    # Case C: Chart
                    elif asset.get("type") == "chart_data":
                        try:
                            chart_data = asset.get("data", {})
                            img_path = os.path.join(assets_dir, f"chart_{i}.png")
                            create_bar_chart(chart_data, img_path)
                            right_col.clear()
                            pic = right_col.insert_picture(img_path)
                        except Exception as e:
                            # Create error image if chart fails
                            img_path = os.path.join(assets_dir, f"chart_error_{i}.png")
                            create_code_snippet_image("Chart Error", img_path)
                            right_col.clear()
                            pic = right_col.insert_picture(img_path)
                            
                    # Case D: Formula
                    elif asset.get("type") == "formula_latex":
                        formula = asset.get("content", "")
                        img_path = os.path.join(assets_dir, f"formula_{i}.png")
                        try:
                            create_latex_image(formula, img_path)
                            right_col.clear()
                            pic = right_col.insert_picture(img_path)
                        except Exception as e:
                            # Create error image if formula fails
                            img_path = os.path.join(assets_dir, f"formula_error_{i}.png")
                            create_code_snippet_image("Formula Error", img_path)
                            right_col.clear()
                            pic = right_col.insert_picture(img_path)
    
    # Save the presentation
    prs.save(output_path)

if __name__ == "__main__":
    main()
