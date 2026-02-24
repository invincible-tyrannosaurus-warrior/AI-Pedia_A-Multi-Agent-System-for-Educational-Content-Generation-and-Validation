import json
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt

def delete_all_slides(prs):
    # Must iterate in reverse to avoid index shifting issues
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

def create_code_snippet_image(code, output_path):
    # Create a white background image with code text
    img = Image.new('RGB', (800, 600), color='white')
    draw = ImageDraw.Draw(img)
    
    # Try to use a monospace font for better code display
    try:
        font = ImageFont.truetype("cour.ttf", 16)
    except:
        font = ImageFont.load_default()
    
    # Draw the code text
    draw.text((10, 10), code, fill='black', font=font)
    
    # Save the image
    img.save(output_path)

def create_latex_image(formula, output_path):
    fig = plt.figure(figsize=(3, 1))
    text = fig.text(0.5, 0.5, f"${formula}$", fontsize=20, ha='center', va='center')
    plt.axis('off')
    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close()

def create_bar_chart(data, output_path):
    labels = data.get('labels', [])
    values = data.get('values', [])
    
    # Create a bar chart
    fig, ax = plt.subplots(figsize=(8, 6))
    bars = ax.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
    
    # Customize the chart
    ax.set_xlabel('Categories')
    ax.set_ylabel('Values')
    ax.set_title('Chart Visualization')
    
    # Add value labels on bars
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.01,
                str(value), ha='center', va='bottom')
    
    # Set background to white
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    
    # Save the chart
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

def create_line_chart(data, output_path):
    labels = data.get('labels', [])
    values = data.get('values', [])
    
    # Create a line chart
    fig, ax = plt.subplots(figsize=(8, 6))
    ax.plot(labels, values, marker='o', color='#1f77b4')
    
    # Customize the chart
    ax.set_xlabel('Categories')
    ax.set_ylabel('Values')
    ax.set_title('Line Chart Visualization')
    
    # Set background to white
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    
    # Save the chart
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

def create_pie_chart(data, output_path):
    labels = data.get('labels', [])
    values = data.get('values', [])
    
    # Create a pie chart
    fig, ax = plt.subplots(figsize=(8, 6))
    wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=90)
    
    # Customize the chart
    ax.set_title('Pie Chart Visualization')
    
    # Set background to white
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    
    # Save the chart
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()

def process_visual_asset(asset, assets_dir, idx):
    asset_type = asset.get('type')
    output_path = None
    
    if asset_type == 'code_snippet':
        code = asset.get('content', '')
        output_path = os.path.join(assets_dir, f"code_{idx}.png")
        create_code_snippet_image(code, output_path)
        
    elif asset_type == 'chart_data':
        chart_type = asset.get('chart_type', 'bar')
        data = asset.get('data', {})
        
        output_path = os.path.join(assets_dir, f"chart_{idx}.png")
        
        try:
            if chart_type == 'bar':
                create_bar_chart(data, output_path)
            elif chart_type == 'line':
                create_line_chart(data, output_path)
            elif chart_type == 'pie':
                create_pie_chart(data, output_path)
            else:
                # Default to bar chart
                create_bar_chart(data, output_path)
        except Exception as e:
            # If chart creation fails, create an error image
            print(f"Chart generation failed: {e}")
            error_img = Image.new('RGB', (400, 200), color='white')
            draw = ImageDraw.Draw(error_img)
            draw.text((10, 10), "Chart Error", fill='red')
            error_img.save(output_path)
            
    elif asset_type == 'formula_latex':
        formula = asset.get('content', '')
        output_path = os.path.join(assets_dir, f"formula_{idx}.png")
        create_latex_image(formula, output_path)
        
    return output_path

# Load storyboard
with open(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_31894469dd264514a5d4bd9e1e5a3d62/storyboard.json", "r", encoding="utf-8") as f:
    storyboard = json.load(f)

slides_data = storyboard.get("slides", [])

# Load presentation template
template_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx"
prs = Presentation(template_path)

# Delete all existing slides
delete_all_slides(prs)

# Define assets directory
assets_dir = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_31894469dd264514a5d4bd9e1e5a3d62/assets"

# Process each slide
for slide_data in slides_data:
    slide_type = slide_data.get('type', 'content')
    
    if slide_type == 'title':
        # Create title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
        
        # Set title and subtitle
        title_placeholder = slide.placeholders[0]
        subtitle_placeholder = slide.placeholders[1]
        
        title_placeholder.text = slide_data.get('title', 'Untitled')
        subtitle_placeholder.text = slide_data.get('subtitle', '')
        
    elif slide_type == 'content':
        content_type = slide_data.get('content_type', 'single_column')
        
        if content_type == 'single_column':
            # Create single column content slide
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Content Slide (1 Column)
            
            # Set title
            title_placeholder = slide.placeholders[0]
            title_placeholder.text = slide_data.get('title', 'Untitled')
            
            # Set body text
            body_placeholder = slide.placeholders[11]
            body_text = '\n'.join(slide_data.get('content', []))
            body_placeholder.text = body_text
            
        elif content_type == 'double_column':
            # Create double column content slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Content Slide (2 Columns)
            
            # Set title
            title_placeholder = slide.placeholders[0]
            title_placeholder.text = slide_data.get('title', 'Untitled')
            
            # Set left column text
            left_placeholder = slide.placeholders[12]
            left_content = '\n'.join(slide_data.get('left_column', []))
            left_placeholder.text = left_content
            
            # Handle visual assets in right column
            visual_assets = slide_data.get('visual_assets', [])
            if visual_assets:
                asset = visual_assets[0]  # Take first visual asset
                asset_path = process_visual_asset(asset, assets_dir, 0)
                
                if asset_path and os.path.exists(asset_path):
                    # Get dimensions of the placeholder
                    right_placeholder = slide.placeholders[13]
                    slide.shapes.add_picture(
                        asset_path,
                        right_placeholder.left,
                        right_placeholder.top,
                        right_placeholder.width,
                        right_placeholder.height
                    )

# Save the presentation
output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_31894469dd264514a5d4bd9e1e5a3d62/lesson_1769359890.pptx"
prs.save(output_path)
