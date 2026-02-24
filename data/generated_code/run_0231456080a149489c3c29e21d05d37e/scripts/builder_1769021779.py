from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import json
import os

def build_presentation():
    # Load storyboard
    with open(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_0231456080a149489c3c29e21d05d37e/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx", 5)
        
        # Add new slide with specified layout
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # Get title
        title = slide_data.get("title", "")
        
        # Set title if placeholder exists
        try:
            slide.placeholders[0].text = title
        except KeyError:
            pass
        
        # Handle different layouts
        if layout_idx == 5:  # 1 Column Content Slide
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[11]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # 2 Column Content Slide
            # Left column - content
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[12]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Right column - visual assets
            visual_assets = slide_data.get("visual_assets", None)
            if visual_assets:
                # Handle case where visual_assets is a list
                if isinstance(visual_assets, list):
                    asset = visual_assets[0] if visual_assets else None
                else:
                    asset = visual_assets
                
                # Handle different types of assets
                if isinstance(asset, str):
                    # It's a path to an image
                    if os.path.exists(asset):
                        try:
                            ph = slide.placeholders[13]
                            slide.shapes.add_picture(asset, ph.left, ph.top, ph.width, ph.height)
                        except KeyError:
                            pass
                elif isinstance(asset, dict) and asset.get("type") == "code_snippet":
                    # Generate code snippet image
                    code = asset.get("code", "")
                    img_path = f"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_0231456080a149489c3c29e21d05d37e/assets/generated_code_{slide_idx}.png"
                    
                    # Create image with PIL
                    img = Image.new('RGB', (800, 600), 'white')
                    draw = ImageDraw.Draw(img)
                    
                    # Try to use a monospace font if available
                    try:
                        font = ImageFont.truetype("cour.ttf", 16)
                    except:
                        font = ImageFont.load_default()
                    
                    # Draw text
                    draw.text((10, 10), code, fill='black', font=font)
                    
                    # Save image
                    img.save(img_path)
                    
                    # Insert into slide
                    try:
                        ph = slide.placeholders[13]
                        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                    except KeyError:
                        pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_0231456080a149489c3c29e21d05d37e/lesson_1769021779.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
