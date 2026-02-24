import json
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt

def create_latex_image(formula, output_path):
    """Create a LaTeX formula image using matplotlib."""
    fig = plt.figure(figsize=(3, 1))
    text = fig.text(0.5, 0.5, f"${formula}$", fontsize=20, ha='center', va='center')
    plt.axis('off')
    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close()

def create_code_snippet_image(code, output_path):
    """Create an image of a code snippet using PIL."""
    # Create a blank image with white background
    img = Image.new('RGB', (600, 400), color='white')
    draw = ImageDraw.Draw(img)
    
    # Load a font (fallback to default if not available)
    try:
        font = ImageFont.truetype("arial.ttf", 16)
    except:
        font = ImageFont.load_default()
    
    # Draw the code text
    draw.text((10, 10), code, fill='black', font=font)
    
    # Save the image
    img.save(output_path)

def create_chart_image(data, output_path, chart_type='bar'):
    """Create a chart image using matplotlib."""
    try:
        labels = data.get('labels', [])
        values = data.get('values', [])
        
        # Set up the figure
        fig, ax = plt.subplots(figsize=(6, 4))
        fig.patch.set_alpha(0)
        ax.patch.set_alpha(0)
        
        # Choose chart type
        if chart_type == 'bar':
            bars = ax.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
        elif chart_type == 'line':
            ax.plot(labels, values, marker='o', color='#1f77b4')
        elif chart_type == 'pie':
            ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=90)
        else:
            bars = ax.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
            
        # Customize chart
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        ax.grid(True, alpha=0.3)
        
        # Adjust layout to prevent clipping
        plt.tight_layout()
        
        # Save the chart
        plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=False, facecolor='white')
        plt.close()
    except Exception as e:
        # In case of error, create a placeholder image
        print(f"Error creating chart: {e}")
        img = Image.new('RGB', (600, 400), color='white')
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 20)
        except:
            font = ImageFont.load_default()
        draw.text((10, 10), "Chart Error", fill='red', font=font)
        img.save(output_path)

def delete_all_slides(prs):
    """Delete all slides from presentation."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for s in slides:
        xml_slides.remove(s)

# Load storyboard
storyboard_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_06e020693af846a298b85e0fa29861cc/storyboard.json"
with open(storyboard_path, "r", encoding="utf-8") as f:
    storyboard = json.load(f)

slides_data = storyboard.get("slides", [])

# Load presentation template
template_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx"
prs = Presentation(template_path)

# Delete all existing slides
delete_all_slides(prs)

# Define asset directory
asset_dir = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/assets"

# Process each slide
for i, slide_data in enumerate(slides_data):
    slide_type = slide_data.get("type", "content")
    title = slide_data.get("title", "")
    
    if slide_type == "title":
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
        title_placeholder = slide.placeholders[0]
        subtitle_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        subtitle_placeholder.text = slide_data.get("subtitle", "")
        
    elif slide_type == "content":
        # Content slide (1 column)
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Content Slide layout
        title_placeholder = slide.placeholders[0]
        body_placeholder = slide.placeholders[11]
        
        title_placeholder.text = title
        body_placeholder.text = slide_data.get("content", "")
        
        # Set text color to black
        for paragraph in body_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)
                
    elif slide_type == "content_2col":
        # Content slide (2 columns)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Two-column layout
        title_placeholder = slide.placeholders[0]
        left_placeholder = slide.placeholders[12]
        right_placeholder = slide.placeholders[13]
        
        title_placeholder.text = title
        
        # Left column content
        left_content = slide_data.get("left_column", "")
        left_placeholder.text = left_content
        
        # Set text color to black for left column
        for paragraph in left_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)
        
        # Handle visual assets in right column
        visual_assets = slide_data.get("visual_assets", [])
        if visual_assets:
            asset = visual_assets[0]  # Take first asset
            
            if isinstance(asset, str):
                # Direct path to image
                if os.path.exists(asset):
                    # Get the dimensions of the placeholder
                    ph = right_placeholder
                    slide.shapes.add_picture(asset, ph.left, ph.top, ph.width, ph.height)
                else:
                    # Create placeholder text if image doesn't exist
                    img = Image.new('RGB', (400, 300), color='white')
                    draw = ImageDraw.Draw(img)
                    try:
                        font = ImageFont.truetype("arial.ttf", 16)
                    except:
                        font = ImageFont.load_default()
                    draw.text((10, 10), "Image Not Found", fill='red', font=font)
                    img_path = os.path.join(asset_dir, f"placeholder_{i}.png")
                    img.save(img_path)
                    ph = right_placeholder
                    slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                    
            elif isinstance(asset, dict):
                # Asset dictionary
                asset_type = asset.get("type", "")
                
                if asset_type == "code_snippet":
                    # Create code snippet image
                    code = asset.get("content", "")
                    img_path = os.path.join(asset_dir, f"code_{i}.png")
                    create_code_snippet_image(code, img_path)
                    ph = right_placeholder
                    slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                    
                elif asset_type == "chart_data":
                    # Create chart image
                    chart_data = asset.get("data", {})
                    chart_type = asset.get("chart_type", "bar")
                    img_path = os.path.join(asset_dir, f"chart_{i}.png")
                    create_chart_image(chart_data, img_path, chart_type)
                    ph = right_placeholder
                    slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                    
                elif asset_type == "formula_latex":
                    # Create LaTeX formula image
                    formula = asset.get("content", "")
                    img_path = os.path.join(asset_dir, f"formula_{i}.png")
                    create_latex_image(formula, img_path)
                    ph = right_placeholder
                    slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)

# Save the presentation
output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_06e020693af846a298b85e0fa29861cc/lesson_1769286010.pptx"
prs.save(output_path)
