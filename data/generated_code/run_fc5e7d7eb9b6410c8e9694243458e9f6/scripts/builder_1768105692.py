from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/archive/data/generated_code/run_fc5e7d7eb9b6410c8e9694243458e9f6/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/archive/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx")
        
        # Add new slide with specified layout
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title if present
        title = slide_data.get("title", "")
        try:
            slide.placeholders[0].text = title
        except KeyError:
            pass
        
        # Handle different layouts
        if layout_idx == 5:  # Content Slide (1 Column)
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[11]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            # Left column - content
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[12]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Right column - visual assets
            visual_assets = slide_data.get("visual_assets", None)
            if visual_assets:
                if isinstance(visual_assets, list):
                    asset = visual_assets[0] if visual_assets else None
                else:
                    asset = visual_assets
                    
                if isinstance(asset, dict) and asset.get("type") == "code_snippet":
                    # Generate code snippet image
                    code = asset.get("code", "")
                    img_path = f"D:/L3/Individual_project/AI_Pedia_Local/archive/data/generated_code/run_fc5e7d7eb9b6410c8e9694243458e9f6/assets/generated_code_{slide_idx}.png"
                    
                    # Create image with PIL
                    img = Image.new('RGB', (800, 600), 'white')
                    draw = ImageDraw.Draw(img)
                    # Simple monospace font drawing
                    font_size = 16
                    y_position = 10
                    lines = code.split('\n')
                    for line in lines:
                        draw.text((10, y_position), line, fill='black')
                        y_position += font_size + 2
                    
                    img.save(img_path)
                    asset_path = img_path
                elif isinstance(asset, str):
                    # Assume it's a path to an existing image
                    asset_path = asset
                else:
                    asset_path = None
                
                # Insert image into placeholder
                if asset_path and os.path.exists(asset_path):
                    try:
                        ph = slide.placeholders[13]
                        slide.shapes.add_picture(asset_path, ph.left, ph.top, ph.width, ph.height)
                    except KeyError:
                        pass

    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/archive/data/slides/lesson_1768105692.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
