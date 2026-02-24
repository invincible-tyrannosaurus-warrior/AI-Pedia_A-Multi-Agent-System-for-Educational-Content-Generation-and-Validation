from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import json
import os

def build_presentation():
    # Load storyboard data
    storyboard_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_2fe8c75ead1140979fdabb058bd1e497/storyboard.json"
    with open(storyboard_path, "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    template_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx"
    prs = Presentation(template_path)
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx")
        
        # Add new slide based on layout index
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            try:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
            except KeyError:
                pass
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            
            # Set body text with black color for all paragraphs
            try:
                body_shape = slide.placeholders[11]
                content = slide_data.get("content", "")
                body_shape.text = content
                for p in body_shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            
            # Set left column content with black color for all paragraphs
            try:
                left_shape = slide.placeholders[12]
                content = slide_data.get("content", "")
                left_shape.text = content
                for p in left_shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Handle visual assets in right column
            try:
                visual_assets = slide_data.get("visual_assets", [])
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset = visual_assets[0]
                    
                    if isinstance(asset, str):  # File path
                        asset_path = asset
                        if os.path.exists(asset_path):
                            ph = slide.placeholders[13]
                            slide.shapes.add_picture(asset_path, ph.left, ph.top, ph.width, ph.height)
                            
                    elif isinstance(asset, dict) and asset.get("type") == "code_snippet":  # Code snippet
                        code = asset.get("code", "")
                        
                        # Generate image with code snippet
                        img = Image.new('RGB', (800, 600), 'white')
                        draw = ImageDraw.Draw(img)
                        
                        # Try to use a monospace font if available
                        try:
                            font = ImageFont.truetype("cour.ttf", 16)
                        except:
                            font = ImageFont.load_default()
                        
                        # Draw code text
                        draw.text((10, 10), code, fill='black', font=font)
                        
                        # Save image to assets directory
                        assets_dir = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/assets"
                        img_path = os.path.join(assets_dir, f"generated_code_{slide_idx}.png")
                        img.save(img_path)
                        
                        # Insert image into placeholder
                        ph = slide.placeholders[13]
                        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                        
            except KeyError:
                pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_2fe8c75ead1140979fdabb058bd1e497/lesson_1769022337.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
