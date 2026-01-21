from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import json
import os

def build_presentation():
    # Load storyboard
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_6226a581d6d3431daab082631ba6f48c/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx")
        
        # Add new slide based on layout
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            try:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
            except KeyError:
                pass
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            
            # Handle body text with black color for all paragraphs
            try:
                content = slide_data.get("content", "")
                slide.placeholders[11].text = content
                for p in slide.placeholders[11].text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            
            # Left column - handle body text
            try:
                content = slide_data.get("content", "")
                slide.placeholders[12].text = content
                for p in slide.placeholders[12].text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Right column - handle visual assets
            try:
                visual_assets = slide_data.get("visual_assets", None)
                if visual_assets:
                    # If it's a list, take the first item
                    if isinstance(visual_assets, list):
                        visual_asset = visual_assets[0]
                    else:
                        visual_asset = visual_assets
                    
                    # Case A: String path to an image file
                    if isinstance(visual_asset, str) and os.path.exists(visual_asset):
                        slide.placeholders[13].insert_picture(visual_asset)
                    
                    # Case B: Dictionary representing a code snippet
                    elif isinstance(visual_asset, dict) and visual_asset.get("type") == "code_snippet":
                        code_content = visual_asset.get("code", "")
                        
                        # Generate image with code snippet
                        img = Image.new('RGB', (800, 600), 'white')
                        draw = ImageDraw.Draw(img)
                        
                        # Try to use a monospace font if available
                        try:
                            font = ImageFont.truetype("cour.ttf", 16)
                        except:
                            font = ImageFont.load_default()
                        
                        # Draw code text
                        draw.text((10, 10), code_content, fill='black', font=font)
                        
                        # Save image to assets directory
                        asset_path = f"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_6226a581d6d3431daab082631ba6f48c/assets/generated_code_{slide_idx}.png"
                        img.save(asset_path)
                        
                        # Insert picture into slide
                        slide.placeholders[13].insert_picture(asset_path)
            except KeyError:
                pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768042492.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
