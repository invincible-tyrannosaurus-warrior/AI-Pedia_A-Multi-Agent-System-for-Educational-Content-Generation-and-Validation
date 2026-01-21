from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_913b1b2feced4e0ab2cf6db2f6587ead/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide in storyboard
    for slide_data in slides_data:
        layout_idx = slide_data.get("layout_idx")
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        visual_assets = slide_data.get("visual_assets", [])
        
        # Create new slide based on layout
        if layout_idx == 0:  # Title Slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if 0 in slide.placeholders:
                title_shape = slide.placeholders[0]
                title_shape.text = title
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Set subtitle
            if 1 in slide.placeholders:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = "Generated Presentation"
                subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if 0 in slide.placeholders:
                title_shape = slide.placeholders[0]
                title_shape.text = title
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Set body text
            if 11 in slide.placeholders:
                body_shape = slide.placeholders[11]
                body_shape.text = content
                body_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if 0 in slide.placeholders:
                title_shape = slide.placeholders[0]
                title_shape.text = title
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Set left column content
            if 12 in slide.placeholders:
                left_shape = slide.placeholders[12]
                left_shape.text = content
                left_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Handle right column visual asset
            if 13 in slide.placeholders and visual_assets:
                # Get first visual asset
                asset_path = visual_assets[0] if isinstance(visual_assets, list) else visual_assets
                
                # If it's a code snippet or chart, we need to generate an image
                if isinstance(asset_path, dict) and 'type' in asset_path:
                    # For now, just create a placeholder image
                    img_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_913b1b2feced4e0ab2cf6db2f6587ead/assets/placeholder.png"
                    
                    # Create a simple image with text
                    img = Image.new('RGB', (400, 300), color='white')
                    draw = ImageDraw.Draw(img)
                    draw.text((10, 10), f"Visual Asset: {asset_path.get('type', 'Unknown')}", fill='black')
                    img.save(img_path)
                    asset_path = img_path
                
                # Insert image into placeholder
                try:
                    slide.placeholders[13].insert_picture(asset_path)
                except Exception as e:
                    print(f"Could not insert image {asset_path}: {e}")
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768027624.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
