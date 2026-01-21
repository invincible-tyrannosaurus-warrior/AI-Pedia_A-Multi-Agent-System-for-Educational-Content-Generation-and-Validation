from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_913b1b2feced4e0ab2cf6db2f6587ead/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide in the storyboard
    for slide_data in slides_data:
        layout_idx = slide_data.get("layout_idx")
        
        # Create new slide based on layout index
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.placeholders[0]
            subtitle = slide.placeholders[1]
            
            if title and 'title' in slide_data:
                title.text = slide_data['title']
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            if subtitle and 'subtitle' in slide_data:
                subtitle.text = slide_data['subtitle']
                subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.placeholders[0]
            body = slide.placeholders[11]
            
            if title and 'title' in slide_data:
                title.text = slide_data['title']
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            if body and 'content' in slide_data:
                body.text = slide_data['content']
                body.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.placeholders[0]
            left_column = slide.placeholders[12]
            right_column = slide.placeholders[13]
            
            if title and 'title' in slide_data:
                title.text = slide_data['title']
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
            if left_column and 'content' in slide_data:
                left_column.text = slide_data['content']
                left_column.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
            if right_column and 'visual_assets' in slide_data:
                visual_assets = slide_data['visual_assets']
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset_path = visual_assets[0]
                    # Check if it's an image path or needs to be generated
                    if asset_path.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                        # It's already an image, just copy it to assets directory
                        img_filename = os.path.basename(asset_path)
                        target_path = os.path.join(
                            r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_913b1b2feced4e0ab2cf6db2f6587ead/assets",
                            img_filename
                        )
                        # Copy or move the file if needed
                        if not os.path.exists(target_path):
                            import shutil
                            shutil.copy(asset_path, target_path)
                        # Insert image into placeholder
                        try:
                            right_column.clear()
                            pic = right_column.insert_picture(target_path)
                        except Exception as e:
                            print(f"Error inserting image {target_path}: {e}")
                    else:
                        # Assume it's code or chart to be generated
                        # For now we'll skip generation but note it in the placeholder
                        right_column.text = "[Visual Asset Not Generated]"
                        right_column.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                else:
                    # No visual asset provided
                    right_column.text = ""
                    right_column.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Save the presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768027587.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

# Run the function with error handling
try:
    build_presentation()
except Exception as e:
    print(f"An error occurred while building the presentation: {e}")
