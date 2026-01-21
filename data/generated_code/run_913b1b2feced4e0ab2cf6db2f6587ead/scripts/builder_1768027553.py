from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_913b1b2feced4e0ab2cf6db2f6587ead/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide in the storyboard
    for slide_data in slides_data:
        layout_idx = slide_data.get("layout_idx")
        
        # Create new slide based on layout index
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.placeholders[0]
            subtitle = slide.placeholders[1]
            
            if title is not None:
                title.text = slide_data.get("title", "")
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
            if subtitle is not None:
                subtitle.text = slide_data.get("subtitle", "")
                subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.placeholders[0]
            body = slide.placeholders[11]
            
            if title is not None:
                title.text = slide_data.get("title", "")
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
            if body is not None:
                content = slide_data.get("content", "")
                body.text = content
                body.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.placeholders[0]
            left_column = slide.placeholders[12]
            right_column = slide.placeholders[13]
            
            # Set title
            if title is not None:
                title.text = slide_data.get("title", "")
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Set left column content
            if left_column is not None:
                content = slide_data.get("content", "")
                left_column.text = content
                left_column.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Handle visual assets in right column
            if right_column is not None:
                visual_assets = slide_data.get("visual_assets", [])
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset_path = visual_assets[0]
                    # If it's a code snippet or chart, we need to generate an image
                    if asset_path.endswith(".png"):
                        # Just copy the image if it already exists
                        pass
                    else:
                        # For now, we'll skip generating images since we don't have generation logic
                        # In a real implementation, you would generate the image here
                        pass
                else:
                    # No visual assets to display
                    pass
    
    # Save the presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768027553.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"An error occurred: {e}")
