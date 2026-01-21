from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_fe4cc96e8e71405582195f61f72e622c/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx")
        
        # Add new slide based on layout
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            try:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
            except KeyError:
                pass
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            
            # Set body text with black color for all paragraphs
            try:
                body_shape = slide.placeholders[11]
                content = slide_data.get("content", "")
                body_shape.text = content
                for p in body_shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            
            # Set left column content
            try:
                left_shape = slide.placeholders[12]
                content = slide_data.get("content", "")
                left_shape.text = content
                for p in left_shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Handle right column visual assets
            try:
                visual_assets = slide_data.get("visual_assets", [])
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset = visual_assets[0]
                    
                    if isinstance(asset, str):
                        # Case A: String path
                        asset_path = asset
                        if os.path.exists(asset_path):
                            ph = slide.placeholders[13]
                            slide.shapes.add_picture(asset_path, ph.left, ph.top, ph.width, ph.height)
                    elif isinstance(asset, dict) and asset.get("type") == "code_snippet":
                        # Case B: Code snippet
                        code = asset.get("code", "")
                        img_path = f"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_fe4cc96e8e71405582195f61f72e622c/assets/generated_code_{slide_idx}.png"
                        
                        # Create image with code
                        img = Image.new('RGB', (800, 600), 'white')
                        draw = ImageDraw.Draw(img)
                        
                        # Try to use a system font
                        try:
                            font = ImageFont.truetype("arial.ttf", 16)
                        except:
                            font = ImageFont.load_default()
                        
                        # Draw code text
                        draw.text((10, 10), code, fill='black', font=font)
                        
                        # Save image
                        img.save(img_path)
                        
                        # Insert into slide
                        ph = slide.placeholders[13]
                        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
            except KeyError:
                pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768043563.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
