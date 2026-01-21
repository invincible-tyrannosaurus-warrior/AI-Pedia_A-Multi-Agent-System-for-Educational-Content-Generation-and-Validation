from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_034a2b6ce8404d299eb3abcf3f33fcbd/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_data in slides_data:
        layout_idx = slide_data.get("layout_idx")
        title = slide_data.get("title", "")
        
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            if 0 in slide.placeholders:
                slide.placeholders[0].text = title
                slide.placeholders[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            if 1 in slide.placeholders:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
                slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            if 0 in slide.placeholders:
                slide.placeholders[0].text = title
                slide.placeholders[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            if 11 in slide.placeholders:
                content = slide_data.get("content", "")
                slide.placeholders[11].text = content
                slide.placeholders[11].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            if 0 in slide.placeholders:
                slide.placeholders[0].text = title
                slide.placeholders[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            if 12 in slide.placeholders:
                content = slide_data.get("content", "")
                slide.placeholders[12].text = content
                slide.placeholders[12].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            if 13 in slide.placeholders:
                visual_assets = slide_data.get("visual_assets", [])
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset_path = visual_assets[0]
                    # Save image to assets directory
                    filename = os.path.basename(asset_path)
                    target_path = os.path.join(
                        r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_034a2b6ce8404d299eb3abcf3f33fcbd/assets",
                        filename
                    )
                    # Copy or process the image as needed
                    # For now, assuming it's already generated at the expected path
                    if os.path.exists(asset_path):
                        # If we need to generate an image instead of just copying
                        # This would be where you'd generate your chart/code snippet
                        pass
                    
                    # Insert image into placeholder
                    left = Inches(0)
                    top = Inches(0)
                    width = Inches(5)
                    height = Inches(3)
                    slide.placeholders[13].insert_picture(target_path)
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768026953.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"An error occurred: {e}")
