from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_034a2b6ce8404d299eb3abcf3f33fcbd/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide in the storyboard
    for slide_data in slides_data:
        layout_idx = slide_data.get("layout_idx")
        title = slide_data.get("title", "")
        
        # Create new slide based on layout index
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            if 0 in slide.placeholders:
                slide.placeholders[0].text = title
                slide.placeholders[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            if 1 in slide.placeholders:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
                slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if 0 in slide.placeholders:
                slide.placeholders[0].text = title
                slide.placeholders[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            if 11 in slide.placeholders:
                content = slide_data.get("content", "")
                slide.placeholders[11].text = content
                slide.placeholders[11].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            if 0 in slide.placeholders:
                slide.placeholders[0].text = title
                slide.placeholders[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
            # Left column - body text
            if 12 in slide.placeholders:
                content = slide_data.get("content", "")
                slide.placeholders[12].text = content
                slide.placeholders[12].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
            # Right column - visual asset
            if 13 in slide.placeholders:
                visual_assets = slide_data.get("visual_assets", [])
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset_path = visual_assets[0]
                    
                    # If it's an image path, just insert it
                    if os.path.exists(asset_path):
                        # Insert existing image
                        left = Inches(6.5)
                        top = Inches(1.5)
                        width = Inches(4)
                        height = Inches(3)
                        slide.shapes.add_picture(asset_path, left, top, width, height)
                    else:
                        # It's likely a code snippet or chart to be generated
                        # Generate a simple image with text
                        img_path = os.path.join(
                            r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_034a2b6ce8404d299eb3abcf3f33fcbd/assets",
                            os.path.basename(asset_path)
                        )
                        
                        # Create a simple image with text
                        img = Image.new('RGB', (600, 400), color='white')
                        draw = ImageDraw.Draw(img)
                        draw.text((10, 10), asset_path, fill='black')
                        
                        # Save image
                        img.save(img_path)
                        
                        # Insert image into slide
                        left = Inches(6.5)
                        top = Inches(1.5)
                        width = Inches(4)
                        height = Inches(3)
                        slide.shapes.add_picture(img_path, left, top, width, height)
    
    # Save the presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768027005.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"An error occurred while building the presentation: {e}")
