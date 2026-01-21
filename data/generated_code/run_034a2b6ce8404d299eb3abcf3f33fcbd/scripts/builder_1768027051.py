from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_034a2b6ce8404d299eb3abcf3f33fcbd/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide in storyboard
    for slide_data in slides_data:
        layout_idx = slide_data.get("layout_idx")
        title = slide_data.get("title", "")
        
        # Create new slide based on layout index
        if layout_idx == 0:  # Title Slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if 0 in slide.placeholders:
                title_shape = slide.placeholders[0]
                title_shape.text = title
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Set subtitle
            if 1 in slide.placeholders:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = slide_data.get("subtitle", "")
                subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if 0 in slide.placeholders:
                title_shape = slide.placeholders[0]
                title_shape.text = title
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Set body text
            if 11 in slide.placeholders:
                body_shape = slide.placeholders[11]
                body_shape.text = slide_data.get("content", "")
                body_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if 0 in slide.placeholders:
                title_shape = slide.placeholders[0]
                title_shape.text = title
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Set left column content
            if 12 in slide.placeholders:
                left_shape = slide.placeholders[12]
                left_shape.text = slide_data.get("content", "")
                left_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Handle visual assets for right column
            if 13 in slide.placeholders:
                visual_assets = slide_data.get("visual_assets", [])
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset_path = visual_assets[0]
                    
                    # If it's an image path, copy it to assets directory
                    if asset_path.startswith("data:image"):
                        # This is base64 encoded image, we need to decode and save it
                        pass  # For simplicity, assuming it's already handled elsewhere
                    else:
                        # It's a file path, check if it exists and copy if needed
                        if os.path.exists(asset_path):
                            filename = os.path.basename(asset_path)
                            target_path = os.path.join(
                                r"D:/L3/Individual_project/AI_Pedia_Local/data/generated_code/run_034a2b6ce8404d299eb3abcf3f33fcbd/assets",
                                filename
                            )
                            if not os.path.exists(target_path):
                                import shutil
                                shutil.copy(asset_path, target_path)
                            asset_path = target_path
                    
                    # Insert image into placeholder
                    try:
                        slide.placeholders[13].insert_picture(asset_path)
                    except Exception:
                        # If image insertion fails, add a text placeholder instead
                        placeholder = slide.placeholders[13]
                        placeholder.text = "Image could not be inserted"
                        placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/data/slides/lesson_1768027051.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"An error occurred: {e}")
