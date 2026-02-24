import json
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt

def create_latex_image(formula, output_path):
    fig = plt.figure(figsize=(3, 1))
    text = fig.text(0.5, 0.5, f"${formula}$", fontsize=20, ha='center', va='center')
    plt.axis('off')
    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close()

def delete_all_slides(prs):
    # Access the XML slide list directly
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for s in slides:
        try:
            xml_slides.remove(s)
        except ValueError:
            pass

def create_code_snippet_image(code, output_path):
    # Create an image with the code snippet
    img = Image.new('RGB', (800, 600), color='white')
    draw = ImageDraw.Draw(img)
    
    # Try to load a font, fallback to default if not available
    try:
        font = ImageFont.truetype("arial.ttf", 16)
    except:
        font = ImageFont.load_default()
    
    # Split code into lines and draw each line
    lines = code.split('\n')
    y_position = 10
    for line in lines:
        draw.text((10, y_position), line, fill='black', font=font)
        y_position += 25
    
    img.save(output_path)

def create_chart_image(chart_data, output_path):
    try:
        labels = chart_data.get('labels', [])
        values = chart_data.get('values', [])
        chart_type = chart_data.get('type', 'bar')
        
        plt.figure(figsize=(8, 6))
        plt.style.use('ggplot')
        
        if chart_type == 'bar':
            bars = plt.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
        elif chart_type == 'line':
            plt.plot(labels, values, marker='o', color='#1f77b4')
        elif chart_type == 'pie':
            plt.pie(values, labels=labels, autopct='%1.1f%%', startangle=90)
        else:
            bars = plt.bar(labels, values, color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'])
            
        plt.title(chart_data.get('title', ''), fontsize=16)
        plt.xticks(rotation=45)
        plt.tight_layout()
        
        plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
    except Exception as e:
        # Create error image if chart creation fails
        img = Image.new('RGB', (400, 200), color='white')
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 16)
        except:
            font = ImageFont.load_default()
        draw.text((10, 10), "Chart Error", fill='red', font=font)
        img.save(output_path)

def main():
    # Load storyboard
    storyboard_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_21f2f78a75024a62a514589a9bf390b7/storyboard.json"
    with open(storyboard_path, "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    
    slides_data = storyboard.get("slides", [])
    
    # Define paths
    assets_dir = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_21f2f78a75024a62a514589a9bf390b7/assets"
    template_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx"
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_21f2f78a75024a62a514589a9bf390b7/lesson_1769286396.pptx"
    
    # Load presentation
    prs = Presentation(template_path)
    delete_all_slides(prs)
    
    # Process slides
    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        visual_assets = slide_data.get("visual_assets", [])
        
        # Add new slide based on type
        if slide_type == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
            title_placeholder = slide.placeholders[0]
            subtitle_placeholder = slide.placeholders[1]
            
            title_placeholder.text = title
            subtitle_placeholder.text = content
            
        elif slide_type == "content_1col":
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Content Slide (1 Column)
            title_placeholder = slide.placeholders[0]
            body_placeholder = slide.placeholders[11]
            
            title_placeholder.text = title
            body_placeholder.text = content
            
        elif slide_type == "content_2col":
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Content Slide (2 Columns)
            title_placeholder = slide.placeholders[0]
            left_placeholder = slide.placeholders[12]
            right_placeholder = slide.placeholders[13]
            
            title_placeholder.text = title
            left_placeholder.text = content
            
            # Handle visual assets
            if visual_assets:
                asset = visual_assets[0]  # Take first asset
                
                if isinstance(asset, str):
                    # It's a path to an existing image
                    if os.path.exists(asset):
                        try:
                            slide.shapes.add_picture(asset, right_placeholder.left, right_placeholder.top, right_placeholder.width, right_placeholder.height)
                        except:
                            pass
                elif isinstance(asset, dict):
                    asset_type = asset.get("type")
                    
                    if asset_type == "code_snippet":
                        code_content = asset.get("content", "")
                        img_path = os.path.join(assets_dir, f"code_snippet_{i}.png")
                        create_code_snippet_image(code_content, img_path)
                        try:
                            slide.shapes.add_picture(img_path, right_placeholder.left, right_placeholder.top, right_placeholder.width, right_placeholder.height)
                        except:
                            pass
                            
                    elif asset_type == "chart_data":
                        chart_data = asset.get("data", {})
                        img_path = os.path.join(assets_dir, f"chart_{i}.png")
                        create_chart_image(chart_data, img_path)
                        try:
                            slide.shapes.add_picture(img_path, right_placeholder.left, right_placeholder.top, right_placeholder.width, right_placeholder.height)
                        except:
                            pass
                            
                    elif asset_type == "formula_latex":
                        formula = asset.get("content", "")
                        img_path = os.path.join(assets_dir, f"formula_{i}.png")
                        create_latex_image(formula, img_path)
                        try:
                            slide.shapes.add_picture(img_path, right_placeholder.left, right_placeholder.top, right_placeholder.width, right_placeholder.height)
                        except:
                            pass
    
    # Save presentation
    prs.save(output_path)

if __name__ == "__main__":
    main()
