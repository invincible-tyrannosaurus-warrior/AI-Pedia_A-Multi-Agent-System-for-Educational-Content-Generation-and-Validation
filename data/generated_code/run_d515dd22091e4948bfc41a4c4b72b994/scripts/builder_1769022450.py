from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import json
import os

def build_presentation():
    # Load storyboard data
    with open(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_d515dd22091e4948bfc41a4c4b72b994/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx")
        
        # Add new slide based on layout
        if layout_idx == 0:  # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            try:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
            except KeyError:
                pass
                
        elif layout_idx == 5:  # Content Slide (1 Column)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            try:
                content = slide_data.get("content", "")
                slide.placeholders[11].text = content
                # Set text color to black for all paragraphs
                for p in slide.placeholders[11].text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            try:
                slide.placeholders[0].text = slide_data.get("title", "")
            except KeyError:
                pass
            
            # Left column - content
            try:
                content = slide_data.get("content", "")
                slide.placeholders[12].text = content
                # Set text color to black for all paragraphs
                for p in slide.placeholders[12].text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
            # Right column - visual assets
            try:
                visual_assets = slide_data.get("visual_assets", [])
                if isinstance(visual_assets, list) and len(visual_assets) > 0:
                    asset = visual_assets[0]
                    
                    if isinstance(asset, str):
                        # Check if file exists
                        if os.path.exists(asset):
                            ph = slide.placeholders[13]
                            slide.shapes.add_picture(asset, ph.left, ph.top, ph.width, ph.height)
                        else:
                            # Handle case where file doesn't exist
                            pass
                    elif isinstance(asset, dict) and asset.get("type") == "code_snippet":
                        # Generate code snippet image
                        code = asset.get("code", "")
                        img_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_d515dd22091e4948bfc41a4c4b72b994/assets/generated_code_{}.png".format(slide_idx)
                        
                        # Create image with code
                        img = Image.new('RGB', (800, 600), 'white')
                        draw = ImageDraw.Draw(img)
                        
                        # Try to use a monospace font if available
                        try:
                            font = ImageFont.truetype("cour.ttf", 16)
                        except:
                            font = ImageFont.load_default()
                            
                        # Draw code text
                        draw.text((10, 10), code, fill='black', font=font)
                        
                        # Save image
                        img.save(img_path)
                        
                        # Insert into slide
                        ph = slide.placeholders[13]
                        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
            except KeyError:
                pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_d515dd22091e4948bfc41a4c4b72b994/lesson_1769022450.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
