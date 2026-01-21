from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import json
import os

def build_presentation():
    # Load storyboard
    with open(r"D:/L3/Individual_project/AI_Pedia_Local/archive/data/generated_code/run_f0cebda557ed48f1920381fb0e10271f/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    slides_data = storyboard.get("slides", [])
    
    # Load template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local/archive/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx", 5)
        
        # Add new slide with specified layout
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title if available
        title = slide_data.get("title", "")
        try:
            slide.placeholders[0].text = title
        except KeyError:
            pass
        
        # Handle different layouts
        if layout_idx == 5:  # Content Slide (1 Column)
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[11]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # Content Slide (2 Columns)
            # Left column
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[12]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Right column - handle visual assets
            visual_assets = slide_data.get("visual_assets", None)
            if visual_assets:
                # If it's a list, take the first item
                if isinstance(visual_assets, list):
                    asset = visual_assets[0]
                else:
                    asset = visual_assets
                    
                try:
                    ph = slide.placeholders[13]
                    
                    # Case A: String path to image
                    if isinstance(asset, str) and os.path.isfile(asset):
                        slide.shapes.add_picture(asset, ph.left, ph.top, ph.width, ph.height)
                        
                    # Case B: Dictionary with code snippet
                    elif isinstance(asset, dict) and asset.get("type") == "code_snippet":
                        code = asset.get("code", "")
                        
                        # Generate image with code snippet
                        img_path = f"D:/L3/Individual_project/AI_Pedia_Local/archive/data/generated_code/run_f0cebda557ed48f1920381fb0e10271f/assets/generated_code_{slide_idx}.png"
                        
                        # Create image
                        img = Image.new('RGB', (800, 600), 'white')
                        draw = ImageDraw.Draw(img)
                        
                        # Try to use a monospace font if available
                        try:
                            font = ImageFont.truetype("cour.ttf", 16)
                        except:
                            font = ImageFont.load_default()
                            
                        # Draw code text
                        draw.text((10, 10), code, fill='black', font=font)
                        
                        # Save image
                        img.save(img_path)
                        
                        # Insert into slide
                        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                        
                except KeyError:
                    pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local/archive/data/slides/lesson_1768104809.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
