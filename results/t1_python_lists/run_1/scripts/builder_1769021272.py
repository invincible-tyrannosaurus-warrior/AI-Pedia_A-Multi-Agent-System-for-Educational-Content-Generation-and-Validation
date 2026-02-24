from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import json
import os

def build_presentation():
    # Load storyboard
    with open(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_016fa81c83384394b967c321c1c4bb8f/storyboard.json", "r", encoding="utf-8") as f:
        storyboard = json.load(f)
    
    slides_data = storyboard.get("slides", [])
    
    # Load presentation template
    prs = Presentation(r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/assets/master_template.pptx")
    
    # Robustly clear existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # Iterate in reverse to avoid index shifting issues
    for i in range(len(slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        # Drop relationship to the slide part
        prs.part.drop_rel(rId)
        # Remove the slide element
        del xml_slides[i]
    
    # Process each slide
    for slide_idx, slide_data in enumerate(slides_data):
        layout_idx = slide_data.get("layout_idx", 5)
        
        # Add new slide with specified layout
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title if available
        title = slide_data.get("title", "")
        try:
            slide.placeholders[0].text = title
        except KeyError:
            pass
        
        # Handle different layouts
        if layout_idx == 5:  # 1 Column Content Slide
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[11]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
                
        elif layout_idx == 6:  # 2 Column Content Slide
            # Left column - content
            content = slide_data.get("content", "")
            try:
                shape = slide.placeholders[12]
                shape.text_frame.text = content
                for p in shape.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(0, 0, 0)
            except KeyError:
                pass
            
            # Right column - visual assets
            visual_assets = slide_data.get("visual_assets", None)
            if visual_assets:
                asset_path = None
                
                # Handle case when visual_assets is a list
                if isinstance(visual_assets, list):
                    if len(visual_assets) > 0:
                        visual_assets = visual_assets[0]
                
                # Handle string path
                if isinstance(visual_assets, str):
                    asset_path = visual_assets
                    
                # Handle dictionary (code snippet)
                elif isinstance(visual_assets, dict) and visual_assets.get("type") == "code_snippet":
                    code = visual_assets.get("code", "")
                    
                    # Generate image with code snippet
                    img = Image.new('RGB', (800, 600), 'white')
                    draw = ImageDraw.Draw(img)
                    
                    # Try to use a monospace font if available
                    try:
                        font = ImageFont.truetype("cour.ttf", 16)
                    except:
                        font = ImageFont.load_default()
                    
                    # Draw code text
                    draw.text((10, 10), code, fill='black', font=font)
                    
                    # Save image
                    img_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_016fa81c83384394b967c321c1c4bb8f/assets/generated_code_{}.png".format(slide_idx)
                    img.save(img_path)
                    asset_path = img_path
                
                # Insert image into placeholder
                if asset_path and os.path.exists(asset_path):
                    try:
                        ph = slide.placeholders[13]
                        slide.shapes.add_picture(asset_path, ph.left, ph.top, ph.width, ph.height)
                    except KeyError:
                        pass
    
    # Save presentation
    output_path = r"D:/L3/Individual_project/AI_Pedia_Local_stream/data/generated_code/run_016fa81c83384394b967c321c1c4bb8f/lesson_1769021272.pptx"
    prs.save(output_path)

try:
    build_presentation()
except Exception as e:
    print(f"Error building presentation: {e}")
